from flask import Flask, render_template, request, redirect, session, send_file, url_for, flash, jsonify
import mysql.connector
from dotenv import load_dotenv
import os
import time
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from datetime import date
import re
import json
import traceback
from functools import wraps
from werkzeug.security import generate_password_hash, check_password_hash
from authlib.integrations.flask_client import OAuth

# ------------------ Load ENV ------------------
load_dotenv()

# ------------------ Initialize Flask ------------------
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "supersecretkey")
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"

DB_ERROR_PAGE = (
    "<div style='font-family:system-ui;max-width:500px;margin:80px auto;text-align:center;"
    "background:#1e1233;color:#fff;padding:48px;border-radius:20px;'>"
    "<h1 style='color:#f87171;margin-bottom:16px;'>Database Offline</h1>"
    "<p style='color:#d1d5db;margin-bottom:24px;'>Database connection failed.</p>"
    "<a href='/' style='color:#a78bfa;'>Try again</a></div>"
)

# ------------------ OpenAI ------------------
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
HF_API_KEY = os.getenv("HF_API_KEY", "")

# ------------------ Google OAuth Setup ------------------
oauth = OAuth(app)
google = oauth.register(
    name="google",
    client_id=os.getenv("GOOGLE_CLIENT_ID", ""),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET", ""),
    server_metadata_url="https://accounts.google.com/.well-known/openid-configuration",
    client_kwargs={"scope": "openid email profile"},
)

# ------------------ MySQL (Railway) ------------------
_db = None
_cursor = None

def get_db():
    global _db, _cursor

    if _db is None or not _db.is_connected():
        try:
            _db = mysql.connector.connect(
                host=os.getenv("MYSQLHOST"),
                port=int(os.getenv("MYSQLPORT", 3306)),
                user=os.getenv("MYSQLUSER"),
                password=os.getenv("MYSQLPASSWORD"),
                database=os.getenv("MYSQLDATABASE"),
            )
            _cursor = _db.cursor()
            return _db, _cursor
        except Exception as e:
            print("[Railway MySQL Error]", e)
            return None, None

    return _db, _cursor


def db_required(f):
    """Decorator: returns friendly error page if MySQL is offline."""
    @wraps(f)
    def wrapper(*args, **kwargs):
        try:
            db, cursor = get_db()
            if db is None:
                return DB_ERROR_PAGE, 503
        except Exception:
            return DB_ERROR_PAGE, 503
        return f(*args, **kwargs)
    return wrapper


# ------------------ Design Tokens ------------------
CHARCOAL = RGBColor(0x1A, 0x1A, 0x2E)
SLATE = RGBColor(0x4A, 0x55, 0x68)
MEDIUM_GRAY = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
CARD_RADIUS = 0.04

DEEP_NAVY = RGBColor(0x0F, 0x2B, 0x46)
ROYAL_BLUE = RGBColor(0x1E, 0x5A, 0xA8)
WARM_CORAL = RGBColor(0xE8, 0x63, 0x4A)
TEAL = RGBColor(0x1A, 0x9E, 0x8F)
SOFT_CREAM = RGBColor(0xFA, 0xF8, 0xF5)
LIGHT_BLUE_TINT = RGBColor(0xED, 0xF3, 0xFA)
LIGHT_CORAL_TINT = RGBColor(0xFD, 0xF0, 0xEE)

# ------------------ Translations ------------------
TRANSLATIONS = {
    "English": {
        "dashboard": "Dashboard", "profile": "Profile", "logout": "Logout",
        "create_title": "Create", "stunning_pres": "Stunning Presentations",
        "in_seconds": "in Seconds",
        "hero_subtitle": "Describe your idea and let AI transform it into a professional presentation",
        "describe_pres": "Describe Your Presentation", "generate": "Generate Presentation",
        "tips_title": "Tips for Better Results",
        "tip1": "Be specific about your topic and target audience",
        "tip2": "Mention if you want data, charts, or statistics",
        "tip3": "Specify the tone: professional, casual, educational",
        "tip4": "Include the number of slides if you have a preference",
        "pres_created": "Presentations Created", "happy_users": "Happy Users",
        "satisfaction": "Satisfaction Rate", "social_title": "My Social Media",
        "copyright": "All rights reserved.",
        "profile_title": "My Profile", "email": "Email", "first_name": "First Name",
        "last_name": "Last Name", "language": "Language", "member_since": "Member Since",
        "welcome": "Welcome",
        "history": "History", "history_title": "Presentation History",
        "no_history": "No presentations yet. Create your first one!",
        "download": "Download", "created_on": "Created on",
        "danger_zone": "Danger Zone", "delete_account": "Delete Account",
        "delete_account_desc": "Once you delete your account, there is no going back. All your data and presentations will be permanently removed.",
        "delete_account_confirm": "Are you sure you want to delete your account? This action cannot be undone.",
        "settings": "Settings",
        "loading_ai": "Generating content with AI...",
        "loading_images": "Creating images...",
        "loading_slides": "Building slides...",
        "loading_saving": "Saving presentation...",
        "loading_done": "Done! Downloading...",
        "loading_generating": "Generating a presentation...",
        "loading_completed": "completed",
    },
    "Russian": {
        "dashboard": "\u041f\u0430\u043d\u0435\u043b\u044c", "profile": "\u041f\u0440\u043e\u0444\u0438\u043b\u044c", "logout": "\u0412\u044b\u0439\u0442\u0438",
        "create_title": "\u0421\u043e\u0437\u0434\u0430\u0439\u0442\u0435", "stunning_pres": "\u041f\u043e\u0442\u0440\u044f\u0441\u0430\u044e\u0449\u0438\u0435 \u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0438",
        "in_seconds": "\u0437\u0430 \u0421\u0435\u043a\u0443\u043d\u0434\u044b",
        "hero_subtitle": "\u041e\u043f\u0438\u0448\u0438\u0442\u0435 \u0441\u0432\u043e\u044e \u0438\u0434\u0435\u044e, \u0438 \u0418\u0418 \u043f\u0440\u0435\u0432\u0440\u0430\u0442\u0438\u0442 \u0435\u0451 \u0432 \u043f\u0440\u043e\u0444\u0435\u0441\u0441\u0438\u043e\u043d\u0430\u043b\u044c\u043d\u0443\u044e \u043f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u044e",
        "describe_pres": "\u041e\u043f\u0438\u0448\u0438\u0442\u0435 \u0412\u0430\u0448\u0443 \u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u044e", "generate": "\u0421\u043e\u0437\u0434\u0430\u0442\u044c \u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u044e",
        "tips_title": "\u0421\u043e\u0432\u0435\u0442\u044b \u0434\u043b\u044f \u041b\u0443\u0447\u0448\u0438\u0445 \u0420\u0435\u0437\u0443\u043b\u044c\u0442\u0430\u0442\u043e\u0432",
        "tip1": "\u0411\u0443\u0434\u044c\u0442\u0435 \u043a\u043e\u043d\u043a\u0440\u0435\u0442\u043d\u044b \u0432 \u043e\u0442\u043d\u043e\u0448\u0435\u043d\u0438\u0438 \u0442\u0435\u043c\u044b \u0438 \u0446\u0435\u043b\u0435\u0432\u043e\u0439 \u0430\u0443\u0434\u0438\u0442\u043e\u0440\u0438\u0438",
        "tip2": "\u0423\u043a\u0430\u0436\u0438\u0442\u0435, \u0435\u0441\u043b\u0438 \u043d\u0443\u0436\u043d\u044b \u0434\u0430\u043d\u043d\u044b\u0435, \u0433\u0440\u0430\u0444\u0438\u043a\u0438 \u0438\u043b\u0438 \u0441\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430",
        "tip3": "\u0423\u043a\u0430\u0436\u0438\u0442\u0435 \u0442\u043e\u043d: \u043f\u0440\u043e\u0444\u0435\u0441\u0441\u0438\u043e\u043d\u0430\u043b\u044c\u043d\u044b\u0439, \u043d\u0435\u0444\u043e\u0440\u043c\u0430\u043b\u044c\u043d\u044b\u0439, \u043e\u0431\u0440\u0430\u0437\u043e\u0432\u0430\u0442\u0435\u043b\u044c\u043d\u044b\u0439",
        "tip4": "\u0423\u043a\u0430\u0436\u0438\u0442\u0435 \u043a\u043e\u043b\u0438\u0447\u0435\u0441\u0442\u0432\u043e \u0441\u043b\u0430\u0439\u0434\u043e\u0432, \u0435\u0441\u043b\u0438 \u0435\u0441\u0442\u044c \u043f\u0440\u0435\u0434\u043f\u043e\u0447\u0442\u0435\u043d\u0438\u0435",
        "pres_created": "\u0421\u043e\u0437\u0434\u0430\u043d\u043e \u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0439", "happy_users": "\u0414\u043e\u0432\u043e\u043b\u044c\u043d\u044b\u0445 \u041f\u043e\u043b\u044c\u0437\u043e\u0432\u0430\u0442\u0435\u043b\u0435\u0439",
        "satisfaction": "\u0423\u0440\u043e\u0432\u0435\u043d\u044c \u0423\u0434\u043e\u0432\u043b\u0435\u0442\u0432\u043e\u0440\u0451\u043d\u043d\u043e\u0441\u0442\u0438", "social_title": "\u041c\u043e\u0438 \u0421\u043e\u0446\u0438\u0430\u043b\u044c\u043d\u044b\u0435 \u0421\u0435\u0442\u0438",
        "copyright": "\u0412\u0441\u0435 \u043f\u0440\u0430\u0432\u0430 \u0437\u0430\u0449\u0438\u0449\u0435\u043d\u044b.",
        "profile_title": "\u041c\u043e\u0439 \u041f\u0440\u043e\u0444\u0438\u043b\u044c", "email": "\u042d\u043b. \u043f\u043e\u0447\u0442\u0430", "first_name": "\u0418\u043c\u044f",
        "last_name": "\u0424\u0430\u043c\u0438\u043b\u0438\u044f", "language": "\u042f\u0437\u044b\u043a", "member_since": "\u0423\u0447\u0430\u0441\u0442\u043d\u0438\u043a \u0441",
        "welcome": "\u0414\u043e\u0431\u0440\u043e \u043f\u043e\u0436\u0430\u043b\u043e\u0432\u0430\u0442\u044c",
        "history": "\u0418\u0441\u0442\u043e\u0440\u0438\u044f", "history_title": "\u0418\u0441\u0442\u043e\u0440\u0438\u044f \u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0439",
        "no_history": "\u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0439 \u043f\u043e\u043a\u0430 \u043d\u0435\u0442. \u0421\u043e\u0437\u0434\u0430\u0439\u0442\u0435 \u0441\u0432\u043e\u044e \u043f\u0435\u0440\u0432\u0443\u044e!",
        "download": "\u0421\u043a\u0430\u0447\u0430\u0442\u044c", "created_on": "\u0421\u043e\u0437\u0434\u0430\u043d\u043e",
        "danger_zone": "\u041e\u043f\u0430\u0441\u043d\u0430\u044f \u0437\u043e\u043d\u0430", "delete_account": "\u0423\u0434\u0430\u043b\u0438\u0442\u044c \u0430\u043a\u043a\u0430\u0443\u043d\u0442",
        "delete_account_desc": "\u041f\u043e\u0441\u043b\u0435 \u0443\u0434\u0430\u043b\u0435\u043d\u0438\u044f \u0430\u043a\u043a\u0430\u0443\u043d\u0442\u0430 \u043f\u0443\u0442\u0438 \u043d\u0430\u0437\u0430\u0434 \u043d\u0435\u0442. \u0412\u0441\u0435 \u0434\u0430\u043d\u043d\u044b\u0435 \u0438 \u043f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0438 \u0431\u0443\u0434\u0443\u0442 \u0443\u0434\u0430\u043b\u0435\u043d\u044b \u043d\u0430\u0432\u0441\u0435\u0433\u0434\u0430.",
        "delete_account_confirm": "\u0412\u044b \u0443\u0432\u0435\u0440\u0435\u043d\u044b, \u0447\u0442\u043e \u0445\u043e\u0442\u0438\u0442\u0435 \u0443\u0434\u0430\u043b\u0438\u0442\u044c \u0430\u043a\u043a\u0430\u0443\u043d\u0442? \u042d\u0442\u043e \u0434\u0435\u0439\u0441\u0442\u0432\u0438\u0435 \u043d\u0435\u043b\u044c\u0437\u044f \u043e\u0442\u043c\u0435\u043d\u0438\u0442\u044c.",
        "settings": "\u041d\u0430\u0441\u0442\u0440\u043e\u0439\u043a\u0438",
        "loading_ai": "\u0413\u0435\u043d\u0435\u0440\u0430\u0446\u0438\u044f \u043a\u043e\u043d\u0442\u0435\u043d\u0442\u0430 \u0441 \u0418\u0418...",
        "loading_images": "\u0421\u043e\u0437\u0434\u0430\u043d\u0438\u0435 \u0438\u0437\u043e\u0431\u0440\u0430\u0436\u0435\u043d\u0438\u0439...",
        "loading_slides": "\u0421\u043e\u0437\u0434\u0430\u043d\u0438\u0435 \u0441\u043b\u0430\u0439\u0434\u043e\u0432...",
        "loading_saving": "\u0421\u043e\u0445\u0440\u0430\u043d\u0435\u043d\u0438\u0435 \u043f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u0438...",
        "loading_done": "\u0413\u043e\u0442\u043e\u0432\u043e! \u0421\u043a\u0430\u0447\u0438\u0432\u0430\u043d\u0438\u0435...",
        "loading_generating": "Создание презентации...",
        "loading_completed": "завершено",
    },
    "Spanish": {
        "dashboard": "Panel", "profile": "Perfil", "logout": "Cerrar sesi\u00f3n",
        "create_title": "Crea", "stunning_pres": "Presentaciones Impresionantes",
        "in_seconds": "en Segundos",
        "hero_subtitle": "Describe tu idea y deja que la IA la transforme en una presentaci\u00f3n profesional",
        "describe_pres": "Describe Tu Presentaci\u00f3n", "generate": "Generar Presentaci\u00f3n",
        "tips_title": "Consejos para Mejores Resultados",
        "tip1": "S\u00e9 espec\u00edfico sobre el tema y el p\u00fablico objetivo",
        "tip2": "Menciona si quieres datos, gr\u00e1ficos o estad\u00edsticas",
        "tip3": "Especifica el tono: profesional, casual, educativo",
        "tip4": "Incluye el n\u00famero de diapositivas si tienes preferencia",
        "pres_created": "Presentaciones Creadas", "happy_users": "Usuarios Felices",
        "satisfaction": "Tasa de Satisfacci\u00f3n", "social_title": "Mis Redes Sociales",
        "copyright": "Todos los derechos reservados.",
        "profile_title": "Mi Perfil", "email": "Correo electr\u00f3nico", "first_name": "Nombre",
        "last_name": "Apellido", "language": "Idioma", "member_since": "Miembro desde",
        "welcome": "Bienvenido",
        "history": "Historial", "history_title": "Historial de Presentaciones",
        "no_history": "\u00a1A\u00fan no hay presentaciones. \u00a1Crea tu primera!",
        "download": "Descargar", "created_on": "Creado el",
        "danger_zone": "Zona de Peligro", "delete_account": "Eliminar Cuenta",
        "delete_account_desc": "Una vez que elimines tu cuenta, no hay vuelta atr\u00e1s. Todos tus datos y presentaciones ser\u00e1n eliminados permanentemente.",
        "delete_account_confirm": "\u00bfEst\u00e1s seguro de que quieres eliminar tu cuenta? Esta acci\u00f3n no se puede deshacer.",
        "settings": "Configuraci\u00f3n",
        "loading_ai": "Generando contenido con IA...",
        "loading_images": "Creando im\u00e1genes...",
        "loading_slides": "Construyendo diapositivas...",
        "loading_saving": "Guardando presentaci\u00f3n...",
        "loading_done": "\u00a1Listo! Descargando...",
        "loading_generating": "Generando una presentación...",
        "loading_completed": "completado",
    },
    "French": {
        "dashboard": "Tableau de bord", "profile": "Profil", "logout": "D\u00e9connexion",
        "create_title": "Cr\u00e9ez", "stunning_pres": "Des Pr\u00e9sentations \u00c9poustouflantes",
        "in_seconds": "en Quelques Secondes",
        "hero_subtitle": "D\u00e9crivez votre id\u00e9e et laissez l\u2019IA la transformer en une pr\u00e9sentation professionnelle",
        "describe_pres": "D\u00e9crivez Votre Pr\u00e9sentation", "generate": "G\u00e9n\u00e9rer la Pr\u00e9sentation",
        "tips_title": "Conseils pour de Meilleurs R\u00e9sultats",
        "tip1": "Soyez pr\u00e9cis sur le sujet et le public cible",
        "tip2": "Mentionnez si vous souhaitez des donn\u00e9es, graphiques ou statistiques",
        "tip3": "Pr\u00e9cisez le ton\u00a0: professionnel, d\u00e9contract\u00e9, \u00e9ducatif",
        "tip4": "Indiquez le nombre de diapositives si vous avez une pr\u00e9f\u00e9rence",
        "pres_created": "Pr\u00e9sentations Cr\u00e9\u00e9es", "happy_users": "Utilisateurs Satisfaits",
        "satisfaction": "Taux de Satisfaction", "social_title": "Mes R\u00e9seaux Sociaux",
        "copyright": "Tous droits r\u00e9serv\u00e9s.",
        "profile_title": "Mon Profil", "email": "E-mail", "first_name": "Pr\u00e9nom",
        "last_name": "Nom", "language": "Langue", "member_since": "Membre depuis",
        "welcome": "Bienvenue",
        "history": "Historique", "history_title": "Historique des Pr\u00e9sentations",
        "no_history": "Pas encore de pr\u00e9sentations. Cr\u00e9ez votre premi\u00e8re\u00a0!",
        "download": "T\u00e9l\u00e9charger", "created_on": "Cr\u00e9\u00e9 le",
        "danger_zone": "Zone Dangereuse", "delete_account": "Supprimer le Compte",
        "delete_account_desc": "Une fois votre compte supprim\u00e9, il n'y a pas de retour en arri\u00e8re. Toutes vos donn\u00e9es et pr\u00e9sentations seront d\u00e9finitivement supprim\u00e9es.",
        "delete_account_confirm": "\u00cates-vous s\u00fbr de vouloir supprimer votre compte ? Cette action est irr\u00e9versible.",
        "settings": "Param\u00e8tres",
        "loading_ai": "G\u00e9n\u00e9ration du contenu par IA...",
        "loading_images": "Cr\u00e9ation des images...",
        "loading_slides": "Construction des diapositives...",
        "loading_saving": "Sauvegarde de la pr\u00e9sentation...",
        "loading_done": "Termin\u00e9 ! T\u00e9l\u00e9chargement...",
        "loading_generating": "Génération d'une présentation...",
        "loading_completed": "terminé",
    },
    "German": {
        "dashboard": "\u00dcbersicht", "profile": "Profil", "logout": "Abmelden",
        "create_title": "Erstellen Sie", "stunning_pres": "Beeindruckende Pr\u00e4sentationen",
        "in_seconds": "in Sekunden",
        "hero_subtitle": "Beschreiben Sie Ihre Idee und lassen Sie KI sie in eine professionelle Pr\u00e4sentation verwandeln",
        "describe_pres": "Beschreiben Sie Ihre Pr\u00e4sentation", "generate": "Pr\u00e4sentation Erstellen",
        "tips_title": "Tipps f\u00fcr Bessere Ergebnisse",
        "tip1": "Seien Sie konkret beim Thema und der Zielgruppe",
        "tip2": "Erw\u00e4hnen Sie, ob Sie Daten, Diagramme oder Statistiken m\u00f6chten",
        "tip3": "Geben Sie den Ton an: professionell, locker, lehrreich",
        "tip4": "Geben Sie die Anzahl der Folien an, falls gew\u00fcnscht",
        "pres_created": "Erstellte Pr\u00e4sentationen", "happy_users": "Zufriedene Nutzer",
        "satisfaction": "Zufriedenheitsrate", "social_title": "Meine Sozialen Medien",
        "copyright": "Alle Rechte vorbehalten.",
        "profile_title": "Mein Profil", "email": "E-Mail", "first_name": "Vorname",
        "last_name": "Nachname", "language": "Sprache", "member_since": "Mitglied seit",
        "welcome": "Willkommen",
        "history": "Verlauf", "history_title": "Pr\u00e4sentationsverlauf",
        "no_history": "Noch keine Pr\u00e4sentationen. Erstellen Sie Ihre erste!",
        "download": "Herunterladen", "created_on": "Erstellt am",
        "danger_zone": "Gefahrenzone", "delete_account": "Konto L\u00f6schen",
        "delete_account_desc": "Sobald Sie Ihr Konto l\u00f6schen, gibt es kein Zur\u00fcck. Alle Ihre Daten und Pr\u00e4sentationen werden dauerhaft entfernt.",
        "delete_account_confirm": "Sind Sie sicher, dass Sie Ihr Konto l\u00f6schen m\u00f6chten? Diese Aktion kann nicht r\u00fcckg\u00e4ngig gemacht werden.",
        "settings": "Einstellungen",
        "loading_ai": "Inhalte mit KI generieren...",
        "loading_images": "Bilder erstellen...",
        "loading_slides": "Folien erstellen...",
        "loading_saving": "Pr\u00e4sentation speichern...",
        "loading_done": "Fertig! Wird heruntergeladen...",
        "loading_generating": "Präsentation wird erstellt...",
        "loading_completed": "abgeschlossen",
    },
    "Italian": {
        "dashboard": "Pannello", "profile": "Profilo", "logout": "Esci",
        "create_title": "Crea", "stunning_pres": "Presentazioni Straordinarie",
        "in_seconds": "in Pochi Secondi",
        "hero_subtitle": "Descrivi la tua idea e lascia che l\u2019IA la trasformi in una presentazione professionale",
        "describe_pres": "Descrivi la Tua Presentazione", "generate": "Genera Presentazione",
        "tips_title": "Consigli per Risultati Migliori",
        "tip1": "Sii specifico sull\u2019argomento e il pubblico di riferimento",
        "tip2": "Indica se desideri dati, grafici o statistiche",
        "tip3": "Specifica il tono: professionale, informale, educativo",
        "tip4": "Includi il numero di diapositive se hai una preferenza",
        "pres_created": "Presentazioni Create", "happy_users": "Utenti Soddisfatti",
        "satisfaction": "Tasso di Soddisfazione", "social_title": "I Miei Social Media",
        "copyright": "Tutti i diritti riservati.",
        "profile_title": "Il Mio Profilo", "email": "E-mail", "first_name": "Nome",
        "last_name": "Cognome", "language": "Lingua", "member_since": "Membro dal",
        "welcome": "Benvenuto",
        "history": "Cronologia", "history_title": "Cronologia Presentazioni",
        "no_history": "Nessuna presentazione ancora. Crea la tua prima!",
        "download": "Scarica", "created_on": "Creato il",
        "danger_zone": "Zona Pericolosa", "delete_account": "Elimina Account",
        "delete_account_desc": "Una volta eliminato il tuo account, non si torna indietro. Tutti i tuoi dati e presentazioni saranno rimossi permanentemente.",
        "delete_account_confirm": "Sei sicuro di voler eliminare il tuo account? Questa azione non pu\u00f2 essere annullata.",
        "settings": "Impostazioni",
        "loading_ai": "Generazione contenuti con IA...",
        "loading_images": "Creazione immagini...",
        "loading_slides": "Costruzione diapositive...",
        "loading_saving": "Salvataggio presentazione...",
        "loading_done": "Fatto! Download in corso...",
        "loading_generating": "Generazione di una presentazione...",
        "loading_completed": "completato",
    },
    "Portuguese": {
        "dashboard": "Painel", "profile": "Perfil", "logout": "Sair",
        "create_title": "Crie", "stunning_pres": "Apresenta\u00e7\u00f5es Incr\u00edveis",
        "in_seconds": "em Segundos",
        "hero_subtitle": "Descreva sua ideia e deixe a IA transform\u00e1-la em uma apresenta\u00e7\u00e3o profissional",
        "describe_pres": "Descreva Sua Apresenta\u00e7\u00e3o", "generate": "Gerar Apresenta\u00e7\u00e3o",
        "tips_title": "Dicas para Melhores Resultados",
        "tip1": "Seja espec\u00edfico sobre o tema e o p\u00fablico-alvo",
        "tip2": "Mencione se deseja dados, gr\u00e1ficos ou estat\u00edsticas",
        "tip3": "Especifique o tom: profissional, casual, educativo",
        "tip4": "Inclua o n\u00famero de slides se tiver prefer\u00eancia",
        "pres_created": "Apresenta\u00e7\u00f5es Criadas", "happy_users": "Usu\u00e1rios Felizes",
        "satisfaction": "Taxa de Satisfa\u00e7\u00e3o", "social_title": "Minhas Redes Sociais",
        "copyright": "Todos os direitos reservados.",
        "profile_title": "Meu Perfil", "email": "E-mail", "first_name": "Nome",
        "last_name": "Sobrenome", "language": "Idioma", "member_since": "Membro desde",
        "welcome": "Bem-vindo",
        "history": "Hist\u00f3rico", "history_title": "Hist\u00f3rico de Apresenta\u00e7\u00f5es",
        "no_history": "Nenhuma apresenta\u00e7\u00e3o ainda. Crie sua primeira!",
        "download": "Baixar", "created_on": "Criado em",
        "danger_zone": "Zona de Perigo", "delete_account": "Excluir Conta",
        "delete_account_desc": "Depois de excluir sua conta, n\u00e3o h\u00e1 volta. Todos os seus dados e apresenta\u00e7\u00f5es ser\u00e3o removidos permanentemente.",
        "delete_account_confirm": "Tem certeza de que deseja excluir sua conta? Esta a\u00e7\u00e3o n\u00e3o pode ser desfeita.",
        "settings": "Configura\u00e7\u00f5es",
        "loading_ai": "Gerando conte\u00fado com IA...",
        "loading_images": "Criando imagens...",
        "loading_slides": "Construindo slides...",
        "loading_saving": "Salvando apresenta\u00e7\u00e3o...",
        "loading_done": "Pronto! Baixando...",
        "loading_generating": "Gerando uma apresentação...",
        "loading_completed": "concluído",
    },
    "Chinese": {
        "dashboard": "\u4eea\u8868\u76d8", "profile": "\u4e2a\u4eba\u8d44\u6599", "logout": "\u9000\u51fa",
        "create_title": "\u521b\u5efa", "stunning_pres": "\u60ca\u8273\u7684\u6f14\u793a\u6587\u7a3f",
        "in_seconds": "\u53ea\u9700\u51e0\u79d2",
        "hero_subtitle": "\u63cf\u8ff0\u60a8\u7684\u60f3\u6cd5\uff0c\u8ba9AI\u5c06\u5176\u8f6c\u5316\u4e3a\u4e13\u4e1a\u6f14\u793a\u6587\u7a3f",
        "describe_pres": "\u63cf\u8ff0\u60a8\u7684\u6f14\u793a\u6587\u7a3f", "generate": "\u751f\u6210\u6f14\u793a\u6587\u7a3f",
        "tips_title": "\u83b7\u5f97\u66f4\u597d\u7ed3\u679c\u7684\u6280\u5de7",
        "tip1": "\u660e\u786e\u8bf4\u660e\u4e3b\u9898\u548c\u76ee\u6807\u53d7\u4f17",
        "tip2": "\u8bf4\u660e\u662f\u5426\u9700\u8981\u6570\u636e\u3001\u56fe\u8868\u6216\u7edf\u8ba1\u4fe1\u606f",
        "tip3": "\u6307\u5b9a\u8bed\u6c14\uff1a\u4e13\u4e1a\u3001\u4f11\u95f2\u3001\u6559\u80b2",
        "tip4": "\u5982\u679c\u6709\u504f\u597d\uff0c\u8bf7\u6307\u5b9a\u5e7b\u706f\u7247\u6570\u91cf",
        "pres_created": "\u5df2\u521b\u5efa\u6f14\u793a\u6587\u7a3f", "happy_users": "\u6ee1\u610f\u7528\u6237",
        "satisfaction": "\u6ee1\u610f\u7387", "social_title": "\u6211\u7684\u793e\u4ea4\u5a92\u4f53",
        "copyright": "\u7248\u6743\u6240\u6709\u3002",
        "profile_title": "\u6211\u7684\u4e2a\u4eba\u8d44\u6599", "email": "\u7535\u5b50\u90ae\u4ef6", "first_name": "\u540d",
        "last_name": "\u59d3", "language": "\u8bed\u8a00", "member_since": "\u6ce8\u518c\u65f6\u95f4",
        "welcome": "\u6b22\u8fce",
        "history": "\u5386\u53f2\u8bb0\u5f55", "history_title": "\u6f14\u793a\u6587\u7a3f\u5386\u53f2",
        "no_history": "\u8fd8\u6ca1\u6709\u6f14\u793a\u6587\u7a3f\u3002\u521b\u5efa\u60a8\u7684\u7b2c\u4e00\u4e2a\u5427\uff01",
        "download": "\u4e0b\u8f7d", "created_on": "\u521b\u5efa\u4e8e",
        "danger_zone": "\u5371\u9669\u533a\u57df", "delete_account": "\u5220\u9664\u8d26\u6237",
        "delete_account_desc": "\u5220\u9664\u8d26\u6237\u540e\u65e0\u6cd5\u6062\u590d\u3002\u60a8\u7684\u6240\u6709\u6570\u636e\u548c\u6f14\u793a\u6587\u7a3f\u5c06\u88ab\u6c38\u4e45\u5220\u9664\u3002",
        "delete_account_confirm": "\u60a8\u786e\u5b9a\u8981\u5220\u9664\u8d26\u6237\u5417\uff1f\u6b64\u64cd\u4f5c\u65e0\u6cd5\u64a4\u9500\u3002",
        "settings": "\u8bbe\u7f6e",
        "loading_ai": "AI\u6b63\u5728\u751f\u6210\u5185\u5bb9...",
        "loading_images": "\u6b63\u5728\u521b\u5efa\u56fe\u7247...",
        "loading_slides": "\u6b63\u5728\u6784\u5efa\u5e7b\u706f\u7247...",
        "loading_saving": "\u6b63\u5728\u4fdd\u5b58\u6f14\u793a\u6587\u7a3f...",
        "loading_done": "\u5b8c\u6210\uff01\u6b63\u5728\u4e0b\u8f7d...",
        "loading_generating": "正在生成演示文稿...",
        "loading_completed": "已完成",
    },
    "Japanese": {
        "dashboard": "\u30c0\u30c3\u30b7\u30e5\u30dc\u30fc\u30c9", "profile": "\u30d7\u30ed\u30d5\u30a3\u30fc\u30eb", "logout": "\u30ed\u30b0\u30a2\u30a6\u30c8",
        "create_title": "\u4f5c\u6210", "stunning_pres": "\u7d20\u6674\u3089\u3057\u3044\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3",
        "in_seconds": "\u6570\u79d2\u3067",
        "hero_subtitle": "\u30a2\u30a4\u30c7\u30a2\u3092\u5165\u529b\u3059\u308b\u3060\u3051\u3067\u3001AI\u304c\u30d7\u30ed\u306e\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u306b\u5909\u63db\u3057\u307e\u3059",
        "describe_pres": "\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u3092\u8aac\u660e", "generate": "\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u3092\u751f\u6210",
        "tips_title": "\u3088\u308a\u826f\u3044\u7d50\u679c\u306e\u305f\u3081\u306e\u30d2\u30f3\u30c8",
        "tip1": "\u30c8\u30d4\u30c3\u30af\u3068\u5bfe\u8c61\u8005\u3092\u5177\u4f53\u7684\u306b\u8a18\u8ff0\u3057\u3066\u304f\u3060\u3055\u3044",
        "tip2": "\u30c7\u30fc\u30bf\u3001\u30b0\u30e9\u30d5\u3001\u7d71\u8a08\u304c\u5fc5\u8981\u306a\u5834\u5408\u306f\u8a18\u8f09\u3057\u3066\u304f\u3060\u3055\u3044",
        "tip3": "\u30c8\u30fc\u30f3\u3092\u6307\u5b9a\uff1a\u30d7\u30ed\u30d5\u30a7\u30c3\u30b7\u30e7\u30ca\u30eb\u3001\u30ab\u30b8\u30e5\u30a2\u30eb\u3001\u6559\u80b2\u7684",
        "tip4": "\u5e0c\u671b\u306e\u30b9\u30e9\u30a4\u30c9\u6570\u304c\u3042\u308c\u3070\u6307\u5b9a\u3057\u3066\u304f\u3060\u3055\u3044",
        "pres_created": "\u4f5c\u6210\u6e08\u307f\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3", "happy_users": "\u6e80\u8db3\u3057\u305f\u30e6\u30fc\u30b6\u30fc",
        "satisfaction": "\u6e80\u8db3\u5ea6", "social_title": "\u79c1\u306e\u30bd\u30fc\u30b7\u30e3\u30eb\u30e1\u30c7\u30a3\u30a2",
        "copyright": "\u5168\u8457\u4f5c\u6a29\u6240\u6709\u3002",
        "profile_title": "\u30de\u30a4\u30d7\u30ed\u30d5\u30a3\u30fc\u30eb", "email": "\u30e1\u30fc\u30eb", "first_name": "\u540d",
        "last_name": "\u59d3", "language": "\u8a00\u8a9e", "member_since": "\u767b\u9332\u65e5",
        "welcome": "\u3088\u3046\u3053\u305d",
        "history": "\u5c65\u6b74", "history_title": "\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u5c65\u6b74",
        "no_history": "\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u306f\u307e\u3060\u3042\u308a\u307e\u305b\u3093\u3002\u6700\u521d\u306e\u4e00\u3064\u3092\u4f5c\u6210\u3057\u307e\u3057\u3087\u3046\uff01",
        "download": "\u30c0\u30a6\u30f3\u30ed\u30fc\u30c9", "created_on": "\u4f5c\u6210\u65e5",
        "danger_zone": "\u5371\u967a\u30be\u30fc\u30f3", "delete_account": "\u30a2\u30ab\u30a6\u30f3\u30c8\u524a\u9664",
        "delete_account_desc": "\u30a2\u30ab\u30a6\u30f3\u30c8\u3092\u524a\u9664\u3059\u308b\u3068\u5143\u306b\u623b\u305b\u307e\u305b\u3093\u3002\u3059\u3079\u3066\u306e\u30c7\u30fc\u30bf\u3068\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u304c\u5b8c\u5168\u306b\u524a\u9664\u3055\u308c\u307e\u3059\u3002",
        "delete_account_confirm": "\u30a2\u30ab\u30a6\u30f3\u30c8\u3092\u524a\u9664\u3057\u3066\u3082\u3088\u308d\u3057\u3044\u3067\u3059\u304b\uff1f\u3053\u306e\u64cd\u4f5c\u306f\u53d6\u308a\u6d88\u305b\u307e\u305b\u3093\u3002",
        "settings": "\u8a2d\u5b9a",
        "loading_ai": "AI\u3067\u30b3\u30f3\u30c6\u30f3\u30c4\u3092\u751f\u6210\u4e2d...",
        "loading_images": "\u753b\u50cf\u3092\u4f5c\u6210\u4e2d...",
        "loading_slides": "\u30b9\u30e9\u30a4\u30c9\u3092\u69cb\u7bc9\u4e2d...",
        "loading_saving": "\u30d7\u30ec\u30bc\u30f3\u30c6\u30fc\u30b7\u30e7\u30f3\u3092\u4fdd\u5b58\u4e2d...",
        "loading_done": "\u5b8c\u4e86\uff01\u30c0\u30a6\u30f3\u30ed\u30fc\u30c9\u4e2d...",
        "loading_generating": "プレゼンテーションを生成中...",
        "loading_completed": "完了",
    },
    "Korean": {
        "dashboard": "\ub300\uc2dc\ubcf4\ub4dc", "profile": "\ud504\ub85c\ud544", "logout": "\ub85c\uadf8\uc544\uc6c3",
        "create_title": "\ub9cc\ub4e4\uae30", "stunning_pres": "\uba4b\uc9c4 \ud504\ub808\uc820\ud14c\uc774\uc158",
        "in_seconds": "\uba87 \ucd08 \ub9cc\uc5d0",
        "hero_subtitle": "\uc544\uc774\ub514\uc5b4\ub97c \uc124\uba85\ud558\uba74 AI\uac00 \uc804\ubb38\uc801\uc778 \ud504\ub808\uc820\ud14c\uc774\uc158\uc73c\ub85c \ubcc0\ud658\ud569\ub2c8\ub2e4",
        "describe_pres": "\ud504\ub808\uc820\ud14c\uc774\uc158 \uc124\uba85", "generate": "\ud504\ub808\uc820\ud14c\uc774\uc158 \uc0dd\uc131",
        "tips_title": "\ub354 \ub098\uc740 \uacb0\uacfc\ub97c \uc704\ud55c \ud301",
        "tip1": "\uc8fc\uc81c\uc640 \ub300\uc0c1 \uccad\uc911\uc5d0 \ub300\ud574 \uad6c\uccb4\uc801\uc73c\ub85c \uc791\uc131\ud558\uc138\uc694",
        "tip2": "\ub370\uc774\ud130, \ucc28\ud2b8 \ub610\ub294 \ud1b5\uacc4\uac00 \ud544\uc694\ud55c\uc9c0 \uba85\uc2dc\ud558\uc138\uc694",
        "tip3": "\ud1a4\uc744 \uc9c0\uc815\ud558\uc138\uc694: \uc804\ubb38\uc801, \uce90\uc8fc\uc5bc, \uad50\uc721\uc801",
        "tip4": "\uc6d0\ud558\ub294 \uc2ac\ub77c\uc774\ub4dc \uc218\uac00 \uc788\uc73c\uba74 \ud3ec\ud568\ud558\uc138\uc694",
        "pres_created": "\uc0dd\uc131\ub41c \ud504\ub808\uc820\ud14c\uc774\uc158", "happy_users": "\ub9cc\uc871\ud55c \uc0ac\uc6a9\uc790",
        "satisfaction": "\ub9cc\uc871\ub3c4", "social_title": "\ub0b4 \uc18c\uc15c \ubbf8\ub514\uc5b4",
        "copyright": "\ubaa8\ub4e0 \uad8c\ub9ac \ubcf4\uc720.",
        "profile_title": "\ub0b4 \ud504\ub85c\ud544", "email": "\uc774\uba54\uc77c", "first_name": "\uc774\ub984",
        "last_name": "\uc131", "language": "\uc5b8\uc5b4", "member_since": "\uac00\uc785\uc77c",
        "welcome": "\ud658\uc601\ud569\ub2c8\ub2e4",
        "history": "\uae30\ub85d", "history_title": "\ud504\ub808\uc820\ud14c\uc774\uc158 \uae30\ub85d",
        "no_history": "\uc544\uc9c1 \ud504\ub808\uc820\ud14c\uc774\uc158\uc774 \uc5c6\uc2b5\ub2c8\ub2e4. \uccab \ubc88\uc9f8\ub97c \ub9cc\ub4e4\uc5b4 \ubcf4\uc138\uc694!",
        "download": "\ub2e4\uc6b4\ub85c\ub4dc", "created_on": "\uc0dd\uc131\uc77c",
        "danger_zone": "\uc704\ud5d8 \uad6c\uc5ed", "delete_account": "\uacc4\uc815 \uc0ad\uc81c",
        "delete_account_desc": "\uacc4\uc815\uc744 \uc0ad\uc81c\ud558\uba74 \ub3cc\uc774\ud0ac \uc218 \uc5c6\uc2b5\ub2c8\ub2e4. \ubaa8\ub4e0 \ub370\uc774\ud130\uc640 \ud504\ub808\uc820\ud14c\uc774\uc158\uc774 \uc601\uad6c\uc801\uc73c\ub85c \uc0ad\uc81c\ub429\ub2c8\ub2e4.",
        "delete_account_confirm": "\uacc4\uc815\uc744 \uc0ad\uc81c\ud558\uc2dc\uaca0\uc2b5\ub2c8\uae4c? \uc774 \uc791\uc5c5\uc740 \ucde8\uc18c\ud560 \uc218 \uc5c6\uc2b5\ub2c8\ub2e4.",
        "settings": "\uc124\uc815",
        "loading_ai": "AI\ub85c \ucf58\ud150\uce20 \uc0dd\uc131 \uc911...",
        "loading_images": "\uc774\ubbf8\uc9c0 \uc0dd\uc131 \uc911...",
        "loading_slides": "\uc2ac\ub77c\uc774\ub4dc \uad6c\uc131 \uc911...",
        "loading_saving": "\ud504\ub808\uc820\ud14c\uc774\uc158 \uc800\uc7a5 \uc911...",
        "loading_done": "\uc644\ub8cc! \ub2e4\uc6b4\ub85c\ub4dc \uc911...",
        "loading_generating": "프레젠테이션 생성 중...",
        "loading_completed": "완료",
    },
    "Arabic": {
        "dashboard": "\u0644\u0648\u062d\u0629 \u0627\u0644\u062a\u062d\u0643\u0645", "profile": "\u0627\u0644\u0645\u0644\u0641 \u0627\u0644\u0634\u062e\u0635\u064a", "logout": "\u062a\u0633\u062c\u064a\u0644 \u0627\u0644\u062e\u0631\u0648\u062c",
        "create_title": "\u0623\u0646\u0634\u0626", "stunning_pres": "\u0639\u0631\u0648\u0636 \u062a\u0642\u062f\u064a\u0645\u064a\u0629 \u0645\u0630\u0647\u0644\u0629",
        "in_seconds": "\u0641\u064a \u062b\u0648\u0627\u0646\u064d",
        "hero_subtitle": "\u0635\u0641 \u0641\u0643\u0631\u062a\u0643 \u0648\u062f\u0639 \u0627\u0644\u0630\u0643\u0627\u0621 \u0627\u0644\u0627\u0635\u0637\u0646\u0627\u0639\u064a \u064a\u062d\u0648\u0644\u0647\u0627 \u0625\u0644\u0649 \u0639\u0631\u0636 \u062a\u0642\u062f\u064a\u0645\u064a \u0627\u062d\u062a\u0631\u0627\u0641\u064a",
        "describe_pres": "\u0635\u0641 \u0639\u0631\u0636\u0643 \u0627\u0644\u062a\u0642\u062f\u064a\u0645\u064a", "generate": "\u0625\u0646\u0634\u0627\u0621 \u0627\u0644\u0639\u0631\u0636 \u0627\u0644\u062a\u0642\u062f\u064a\u0645\u064a",
        "tips_title": "\u0646\u0635\u0627\u0626\u062d \u0644\u0646\u062a\u0627\u0626\u062c \u0623\u0641\u0636\u0644",
        "tip1": "\u0643\u0646 \u0645\u062d\u062f\u062f\u064b\u0627 \u0628\u0634\u0623\u0646 \u0627\u0644\u0645\u0648\u0636\u0648\u0639 \u0648\u0627\u0644\u062c\u0645\u0647\u0648\u0631 \u0627\u0644\u0645\u0633\u062a\u0647\u062f\u0641",
        "tip2": "\u0627\u0630\u0643\u0631 \u0625\u0630\u0627 \u0643\u0646\u062a \u062a\u0631\u064a\u062f \u0628\u064a\u0627\u0646\u0627\u062a \u0623\u0648 \u0631\u0633\u0648\u0645 \u0628\u064a\u0627\u0646\u064a\u0629 \u0623\u0648 \u0625\u062d\u0635\u0627\u0626\u064a\u0627\u062a",
        "tip3": "\u062d\u062f\u062f \u0627\u0644\u0646\u0628\u0631\u0629: \u0645\u0647\u0646\u064a\u0629\u060c \u063a\u064a\u0631 \u0631\u0633\u0645\u064a\u0629\u060c \u062a\u0639\u0644\u064a\u0645\u064a\u0629",
        "tip4": "\u062d\u062f\u062f \u0639\u062f\u062f \u0627\u0644\u0634\u0631\u0627\u0626\u062d \u0625\u0630\u0627 \u0643\u0627\u0646 \u0644\u062f\u064a\u0643 \u062a\u0641\u0636\u064a\u0644",
        "pres_created": "\u0627\u0644\u0639\u0631\u0648\u0636 \u0627\u0644\u0645\u064f\u0646\u0634\u0623\u0629", "happy_users": "\u0645\u0633\u062a\u062e\u062f\u0645\u0648\u0646 \u0633\u0639\u062f\u0627\u0621",
        "satisfaction": "\u0645\u0639\u062f\u0644 \u0627\u0644\u0631\u0636\u0627", "social_title": "\u0648\u0633\u0627\u0626\u0644 \u0627\u0644\u062a\u0648\u0627\u0635\u0644 \u0627\u0644\u0627\u062c\u062a\u0645\u0627\u0639\u064a",
        "copyright": "\u062c\u0645\u064a\u0639 \u0627\u0644\u062d\u0642\u0648\u0642 \u0645\u062d\u0641\u0648\u0638\u0629.",
        "profile_title": "\u0645\u0644\u0641\u064a \u0627\u0644\u0634\u062e\u0635\u064a", "email": "\u0627\u0644\u0628\u0631\u064a\u062f \u0627\u0644\u0625\u0644\u0643\u062a\u0631\u0648\u0646\u064a", "first_name": "\u0627\u0644\u0627\u0633\u0645 \u0627\u0644\u0623\u0648\u0644",
        "last_name": "\u0627\u0644\u0643\u0646\u064a\u0629", "language": "\u0627\u0644\u0644\u063a\u0629", "member_since": "\u0639\u0636\u0648 \u0645\u0646\u0630",
        "welcome": "\u0645\u0631\u062d\u0628\u064b\u0627",
        "history": "\u0627\u0644\u0633\u062c\u0644", "history_title": "\u0633\u062c\u0644 \u0627\u0644\u0639\u0631\u0648\u0636 \u0627\u0644\u062a\u0642\u062f\u064a\u0645\u064a\u0629",
        "no_history": "\u0644\u0627 \u062a\u0648\u062c\u062f \u0639\u0631\u0648\u0636 \u062a\u0642\u062f\u064a\u0645\u064a\u0629 \u0628\u0639\u062f. \u0623\u0646\u0634\u0626 \u0623\u0648\u0644 \u0639\u0631\u0636 \u0644\u0643!",
        "download": "\u062a\u062d\u0645\u064a\u0644", "created_on": "\u062a\u0627\u0631\u064a\u062e \u0627\u0644\u0625\u0646\u0634\u0627\u0621",
        "danger_zone": "\u0645\u0646\u0637\u0642\u0629 \u0627\u0644\u062e\u0637\u0631", "delete_account": "\u062d\u0630\u0641 \u0627\u0644\u062d\u0633\u0627\u0628",
        "delete_account_desc": "\u0628\u0645\u062c\u0631\u062f \u062d\u0630\u0641 \u062d\u0633\u0627\u0628\u0643\u060c \u0644\u0627 \u064a\u0645\u0643\u0646 \u0627\u0644\u0639\u0648\u062f\u0629. \u0633\u064a\u062a\u0645 \u0625\u0632\u0627\u0644\u0629 \u062c\u0645\u064a\u0639 \u0628\u064a\u0627\u0646\u0627\u062a\u0643 \u0648\u0639\u0631\u0648\u0636\u0643 \u0627\u0644\u062a\u0642\u062f\u064a\u0645\u064a\u0629 \u0646\u0647\u0627\u0626\u064a\u064b\u0627.",
        "delete_account_confirm": "\u0647\u0644 \u0623\u0646\u062a \u0645\u062a\u0623\u0643\u062f \u0645\u0646 \u0631\u063a\u0628\u062a\u0643 \u0641\u064a \u062d\u0630\u0641 \u062d\u0633\u0627\u0628\u0643\u061f \u0644\u0627 \u064a\u0645\u0643\u0646 \u0627\u0644\u062a\u0631\u0627\u062c\u0639 \u0639\u0646 \u0647\u0630\u0627 \u0627\u0644\u0625\u062c\u0631\u0627\u0621.",
        "settings": "\u0627\u0644\u0625\u0639\u062f\u0627\u062f\u0627\u062a",
        "loading_ai": "...\u062c\u0627\u0631\u064a \u0625\u0646\u0634\u0627\u0621 \u0627\u0644\u0645\u062d\u062a\u0648\u0649 \u0628\u0627\u0644\u0630\u0643\u0627\u0621 \u0627\u0644\u0627\u0635\u0637\u0646\u0627\u0639\u064a",
        "loading_images": "...\u062c\u0627\u0631\u064a \u0625\u0646\u0634\u0627\u0621 \u0627\u0644\u0635\u0648\u0631",
        "loading_slides": "...\u062c\u0627\u0631\u064a \u0628\u0646\u0627\u0621 \u0627\u0644\u0634\u0631\u0627\u0626\u062d",
        "loading_saving": "...\u062c\u0627\u0631\u064a \u062d\u0641\u0638 \u0627\u0644\u0639\u0631\u0636",
        "loading_done": "!\u062a\u0645! \u062c\u0627\u0631\u064a \u0627\u0644\u062a\u0646\u0632\u064a\u0644...",
        "loading_generating": "...جاري إنشاء العرض التقديمي",
        "loading_completed": "مكتمل",
    },
    "Hindi": {
        "dashboard": "\u0921\u0948\u0936\u092c\u094b\u0930\u094d\u0921", "profile": "\u092a\u094d\u0930\u094b\u092b\u093e\u0907\u0932", "logout": "\u0932\u0949\u0917 \u0906\u0909\u091f",
        "create_title": "\u092c\u0928\u093e\u090f\u0901", "stunning_pres": "\u0936\u093e\u0928\u0926\u093e\u0930 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928",
        "in_seconds": "\u0915\u0941\u091b \u0939\u0940 \u0938\u0915\u0947\u0902\u0921 \u092e\u0947\u0902",
        "hero_subtitle": "\u0905\u092a\u0928\u0947 \u0935\u093f\u091a\u093e\u0930 \u0915\u093e \u0935\u0930\u094d\u0923\u0928 \u0915\u0930\u0947\u0902 \u0914\u0930 AI \u0907\u0938\u0947 \u092a\u0947\u0936\u0947\u0935\u0930 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u092e\u0947\u0902 \u092c\u0926\u0932 \u0926\u0947\u0917\u093e",
        "describe_pres": "\u0905\u092a\u0928\u0940 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u0915\u093e \u0935\u0930\u094d\u0923\u0928 \u0915\u0930\u0947\u0902", "generate": "\u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u092c\u0928\u093e\u090f\u0901",
        "tips_title": "\u092c\u0947\u0939\u0924\u0930 \u092a\u0930\u093f\u0923\u093e\u092e\u094b\u0902 \u0915\u0947 \u0932\u093f\u090f \u0938\u0941\u091d\u093e\u0935",
        "tip1": "\u0935\u093f\u0937\u092f \u0914\u0930 \u0932\u0915\u094d\u0937\u094d\u092f \u0926\u0930\u094d\u0936\u0915\u094b\u0902 \u0915\u0947 \u092c\u093e\u0930\u0947 \u092e\u0947\u0902 \u0938\u094d\u092a\u0937\u094d\u091f \u0930\u0939\u0947\u0902",
        "tip2": "\u092c\u0924\u093e\u090f\u0901 \u0915\u094d\u092f\u093e \u0906\u092a\u0915\u094b \u0921\u0947\u091f\u093e, \u091a\u093e\u0930\u094d\u091f \u092f\u093e \u0906\u0901\u0915\u0921\u093c\u0947 \u091a\u093e\u0939\u093f\u090f",
        "tip3": "\u091f\u094b\u0928 \u092c\u0924\u093e\u090f\u0901: \u092a\u0947\u0936\u0947\u0935\u0930, \u0905\u0928\u094c\u092a\u091a\u093e\u0930\u093f\u0915, \u0936\u0948\u0915\u094d\u0937\u093f\u0915",
        "tip4": "\u092f\u0926\u093f \u092a\u0938\u0902\u0926 \u0939\u094b \u0924\u094b \u0938\u094d\u0932\u093e\u0907\u0921 \u0915\u0940 \u0938\u0902\u0916\u094d\u092f\u093e \u0936\u093e\u092e\u093f\u0932 \u0915\u0930\u0947\u0902",
        "pres_created": "\u092c\u0928\u093e\u0908 \u0917\u0908 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928", "happy_users": "\u0916\u0941\u0936 \u0909\u092a\u092f\u094b\u0917\u0915\u0930\u094d\u0924\u093e",
        "satisfaction": "\u0938\u0902\u0924\u0941\u0937\u094d\u091f\u093f \u0926\u0930", "social_title": "\u092e\u0947\u0930\u0947 \u0938\u094b\u0936\u0932 \u092e\u0940\u0921\u093f\u092f\u093e",
        "copyright": "\u0938\u0930\u094d\u0935\u093e\u0927\u093f\u0915\u093e\u0930 \u0938\u0941\u0930\u0915\u094d\u0937\u093f\u0924\u0964",
        "profile_title": "\u092e\u0947\u0930\u0940 \u092a\u094d\u0930\u094b\u092b\u093e\u0907\u0932", "email": "\u0908\u092e\u0947\u0932", "first_name": "\u092a\u0939\u0932\u093e \u0928\u093e\u092e",
        "last_name": "\u0909\u092a\u0928\u093e\u092e", "language": "\u092d\u093e\u0937\u093e", "member_since": "\u0938\u0926\u0938\u094d\u092f \u0924\u093f\u0925\u093f",
        "welcome": "\u0938\u094d\u0935\u093e\u0917\u0924",
        "history": "\u0907\u0924\u093f\u0939\u093e\u0938", "history_title": "\u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u0907\u0924\u093f\u0939\u093e\u0938",
        "no_history": "\u0905\u092d\u0940 \u0915\u094b\u0908 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u0928\u0939\u0940\u0902\u0964 \u0905\u092a\u0928\u093e \u092a\u0939\u0932\u093e \u092c\u0928\u093e\u090f\u0901!",
        "download": "\u0921\u093e\u0909\u0928\u0932\u094b\u0921", "created_on": "\u092c\u0928\u093e\u092f\u093e \u0917\u092f\u093e",
        "danger_zone": "\u0916\u0924\u0930\u0928\u093e\u0915 \u0915\u094d\u0937\u0947\u0924\u094d\u0930", "delete_account": "\u0916\u093e\u0924\u093e \u0939\u091f\u093e\u090f\u0902",
        "delete_account_desc": "\u0916\u093e\u0924\u093e \u0939\u091f\u093e\u0928\u0947 \u0915\u0947 \u092c\u093e\u0926 \u0935\u093e\u092a\u0938 \u0928\u0939\u0940\u0902 \u0906 \u0938\u0915\u0924\u0947\u0964 \u0906\u092a\u0915\u093e \u0938\u092d\u0940 \u0921\u0947\u091f\u093e \u0914\u0930 \u092a\u094d\u0930\u0947\u0938\u0947\u0902\u091f\u0947\u0936\u0928 \u0938\u094d\u0925\u093e\u092f\u0940 \u0930\u0942\u092a \u0938\u0947 \u0939\u091f\u093e \u0926\u093f\u090f \u091c\u093e\u090f\u0902\u0917\u0947\u0964",
        "delete_account_confirm": "\u0915\u094d\u092f\u093e \u0906\u092a \u0935\u093e\u0915\u0908 \u0905\u092a\u0928\u093e \u0916\u093e\u0924\u093e \u0939\u091f\u093e\u0928\u093e \u091a\u093e\u0939\u0924\u0947 \u0939\u0948\u0902? \u092f\u0939 \u0915\u094d\u0930\u093f\u092f\u093e \u0935\u093e\u092a\u0938 \u0928\u0939\u0940\u0902 \u0939\u094b \u0938\u0915\u0924\u0940\u0964",
        "settings": "\u0938\u0947\u091f\u093f\u0902\u0917\u094d\u0938",
        "loading_ai": "AI \u0938\u0947 \u0938\u093e\u092e\u0917\u094d\u0930\u0940 \u092c\u0928\u093e \u0930\u0939\u0947 \u0939\u0948\u0902...",
        "loading_images": "\u091a\u093f\u0924\u094d\u0930 \u092c\u0928\u093e \u0930\u0939\u0947 \u0939\u0948\u0902...",
        "loading_slides": "\u0938\u094d\u0932\u093e\u0907\u0921 \u092c\u0928\u093e \u0930\u0939\u0947 \u0939\u0948\u0902...",
        "loading_saving": "\u092a\u094d\u0930\u0938\u094d\u0924\u0941\u0924\u093f \u0938\u0939\u0947\u091c \u0930\u0939\u0947 \u0939\u0948\u0902...",
        "loading_done": "\u0939\u094b \u0917\u092f\u093e! \u0921\u093e\u0909\u0928\u0932\u094b\u0921 \u0939\u094b \u0930\u0939\u093e \u0939\u0948...",
        "loading_generating": "प्रस्तुति बनाई जा रही है...",
        "loading_completed": "पूर्ण",
    },
    "Uzbek": {
        "dashboard": "Boshqaruv paneli", "profile": "Profil", "logout": "Chiqish",
        "create_title": "Yarating", "stunning_pres": "Ajoyib Taqdimotlar",
        "in_seconds": "Soniyalarda",
        "hero_subtitle": "G\u2018oyangizni tasvirlab bering va sun\u2018iy intellekt uni professional taqdimotga aylantirsin",
        "describe_pres": "Taqdimotingizni Tasvirlang", "generate": "Taqdimot Yaratish",
        "tips_title": "Yaxshi Natijalar Uchun Maslahatlar",
        "tip1": "Mavzu va maqsadli auditoriya haqida aniq bo\u2018ling",
        "tip2": "Ma\u2018lumotlar, grafiklar yoki statistika kerakligini ko\u2018rsating",
        "tip3": "Ohangni belgilang: professional, norasmiy, ta\u2018limiy",
        "tip4": "Agar xohlasangiz slaydlar sonini ko\u2018rsating",
        "pres_created": "Yaratilgan Taqdimotlar", "happy_users": "Mamnun Foydalanuvchilar",
        "satisfaction": "Qoniqish Darajasi", "social_title": "Mening Ijtimoiy Tarmoqlarim",
        "copyright": "Barcha huquqlar himoyalangan.",
        "profile_title": "Mening Profilim", "email": "Elektron pochta", "first_name": "Ism",
        "last_name": "Familiya", "language": "Til", "member_since": "A\u2018zo bo\u2018lgan sana",
        "welcome": "Xush kelibsiz",
        "history": "Tarix", "history_title": "Taqdimotlar Tarixi",
        "no_history": "Hali taqdimotlar yo\u2018q. Birinchisini yarating!",
        "download": "Yuklab olish", "created_on": "Yaratilgan sana",
        "danger_zone": "Xavfli Zona", "delete_account": "Hisobni O'chirish",
        "delete_account_desc": "Hisobingizni o'chirganingizdan so'ng, qaytib bo'lmaydi. Barcha ma'lumotlaringiz va taqdimotlaringiz butunlay o'chiriladi.",
        "delete_account_confirm": "Hisobingizni o'chirishga ishonchingiz komilmi? Bu amalni bekor qilib bo'lmaydi.",
        "settings": "Sozlamalar",
        "loading_ai": "AI bilan kontent yaratilmoqda...",
        "loading_images": "Rasmlar yaratilmoqda...",
        "loading_slides": "Slaydlar qurilmoqda...",
        "loading_saving": "Taqdimot saqlanmoqda...",
        "loading_done": "Tayyor! Yuklab olinmoqda...",
        "loading_generating": "Taqdimot yaratilmoqda...",
        "loading_completed": "bajarildi",
    },
}

LANG_CODES = {
    "English": "en", "Spanish": "es", "French": "fr", "German": "de",
    "Italian": "it", "Portuguese": "pt", "Russian": "ru", "Chinese": "zh",
    "Japanese": "ja", "Korean": "ko", "Arabic": "ar", "Hindi": "hi",
    "Uzbek": "uz",
}


def get_translations():
    lang = session.get("language", "English")
    return TRANSLATIONS.get(lang, TRANSLATIONS["English"])

TOTAL_SLIDES = 10

# Layout grid – standard 16:9 widescreen (13.333" x 7.5")
SW = 13.333
SH = 7.5
MG = 0.60
CT = 1.80
CW = SW - 2 * MG       # ~12.133
GAP = 0.30
IMG_SZ = 5.0

# Derived positions
IMG_RX = SW - MG - IMG_SZ   # ~7.733 (right-side image x)
IMG_TOP = CT                 # 1.80
TXT_BW = IMG_RX - MG - GAP  # ~6.833 (text width beside image)

# Side image dimensions (full-height portrait panel)
SIDE_W = 5.0
SIDE_H = SH  # full slide height

# ------------------ Color Themes ------------------
# 4 curated palettes from user reference
THEMES = {
    "indigo": {
        # Velvet Indigo #4B0082 + near-white warm
        "dark": (0x2D, 0x00, 0x4F), "primary": (0x4B, 0x00, 0x82),
        "secondary": (0x8B, 0x5C, 0xF6), "tertiary": (0x6B, 0x21, 0xA8),
        "light_bg": (0xF8, 0xF7, 0xFC), "card_bg": (0xFF, 0xFF, 0xFF),
        "card_alt": (0xF3, 0xF0, 0xFA),
    },
    "nature": {
        # Green + warm off-white
        "dark": (0x2D, 0x5A, 0x3A), "primary": (0x4A, 0x7C, 0x59),
        "secondary": (0x8B, 0x69, 0x14), "tertiary": (0x5A, 0x9E, 0x6F),
        "light_bg": (0xFB, 0xF9, 0xF5), "card_bg": (0xFF, 0xFF, 0xFF),
        "card_alt": (0xF5, 0xF2, 0xEB),
    },
    "moss": {
        # Moss + cool off-white
        "dark": (0x3A, 0x43, 0x18), "primary": (0x56, 0x61, 0x29),
        "secondary": (0x8B, 0x73, 0x55), "tertiary": (0x7A, 0x8B, 0x45),
        "light_bg": (0xF8, 0xF9, 0xF5), "card_bg": (0xFF, 0xFF, 0xFF),
        "card_alt": (0xF2, 0xF3, 0xEE),
    },
    "navy": {
        # Rich Blue + cool off-white
        "dark": (0x0F, 0x0A, 0x3A), "primary": (0x1A, 0x12, 0x64),
        "secondary": (0x4A, 0x45, 0xB0), "tertiary": (0x2E, 0x28, 0x90),
        "light_bg": (0xF7, 0xF7, 0xFC), "card_bg": (0xFF, 0xFF, 0xFF),
        "card_alt": (0xF0, 0xF0, 0xF8),
    },
}


def apply_theme(name):
    global DEEP_NAVY, ROYAL_BLUE, WARM_CORAL, TEAL
    global SOFT_CREAM, LIGHT_BLUE_TINT, LIGHT_CORAL_TINT
    t = THEMES.get(name, THEMES["indigo"])
    DEEP_NAVY = RGBColor(*t["dark"])
    ROYAL_BLUE = RGBColor(*t["primary"])
    WARM_CORAL = RGBColor(*t["secondary"])
    TEAL = RGBColor(*t["tertiary"])
    SOFT_CREAM = RGBColor(*t["light_bg"])
    LIGHT_BLUE_TINT = RGBColor(*t["card_bg"])
    LIGHT_CORAL_TINT = RGBColor(*t["card_alt"])


# ------------------ Helpers ------------------
def optimize_title(prompt):
    try:
        r = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Return a short presentation title."},
                {"role": "user", "content": prompt}
            ]
        )
        return r.choices[0].message.content.strip()
    except Exception:
        return prompt


def safe_filename(text):
    text = text.replace("\n", " ").replace("\r", " ")
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"\s+", "_", text)
    return text[:50] or "presentation"


_CONTENT_SYSTEM_PROMPT = """You are a presentation content writer. Given a topic, generate content for a 10-slide presentation.
Return ONLY valid JSON (no markdown code fences) with this exact structure:

{
  "title": "Main Presentation Title",
  "subtitle": "Short subtitle or tagline",
  "title_image_prompt": "A vivid illustration or photograph for the title slide ...",
  "theme": "navy",
  "slide_designs": {
    "2": "A", "3": "A", "4": "A", "5": "A", "6": "A", "7": "A", "8": "A", "9": "A"
  },
  "slides": [
    {
      "slide": 2,
      "title": "Slide Title",
      "cards": [
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]}
      ]
    },
    {
      "slide": 3,
      "title": "Slide Title",
      "cards": [
        {"title": "Card Title", "bullets": ["Point 1", "Point 2", "Point 3"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2", "Point 3"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2", "Point 3"]}
      ]
    },
    {
      "slide": 4,
      "title": "Slide Title",
      "image_prompt": "A professional photograph of ...",
      "bullets": ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5"]
    },
    {
      "slide": 5,
      "title": "Slide Title",
      "cards": [
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]},
        {"title": "Card Title", "bullets": ["Point 1", "Point 2"]}
      ]
    },
    {
      "slide": 6,
      "title": "Slide Title",
      "image_prompt": "A professional photograph of ...",
      "bullets": ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5"],
      "stat_number": "85%",
      "stat_label": "of something important"
    },
    {
      "slide": 7,
      "title": "Slide Title",
      "left": {
        "title": "Left Side Title",
        "bullets": ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5"]
      },
      "right": {
        "title": "Right Side Title",
        "bullets": ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5"]
      }
    },
    {
      "slide": 8,
      "title": "Slide Title",
      "steps": [
        {"title": "Step 1", "bullets": ["Detail 1", "Detail 2"]},
        {"title": "Step 2", "bullets": ["Detail 1", "Detail 2"]},
        {"title": "Step 3", "bullets": ["Detail 1", "Detail 2"]},
        {"title": "Step 4", "bullets": ["Detail 1", "Detail 2"]}
      ]
    },
    {
      "slide": 9,
      "title": "Slide Title",
      "image_prompt": "A professional photograph of ...",
      "stat_number": "50%",
      "stat_label": "increase in something",
      "bullets": ["Point 1", "Point 2", "Point 3", "Point 4"]
    },
    {
      "slide": 10,
      "title": "Call-to-Action Title",
      "cards": [
        {"title": "Action 1", "bullets": ["Step 1", "Step 2", "Step 3"]},
        {"title": "Action 2", "bullets": ["Step 1", "Step 2", "Step 3"]},
        {"title": "Action 3", "bullets": ["Step 1", "Step 2", "Step 3"]}
      ]
    }
  ]
}

Slide Design Options (slide_designs):
For slides 2-9, choose "A" or "B". Mix them for visual variety.
- Slide 2: A = full-width dot-badge rows (heading + bullets per row) | B = full-width left-accent cards (provide 2 or 3 cards)
- Slide 3: A = 3 cards in a row | B = 4 cards in a row (provide 3 or 4 cards)
- Slide 4: A = image left + bullets right | B = image right + bullets left
- Slide 5: A = 2x2 card grid with numbered badges | B = big stat columns (card titles = stats: "50M", "98%")
- Slide 6: A = image left + bullets right + stat banner | B = image right + stacked cards left + stat banner
- Slide 7: A = two comparison cards (tinted fills) | B = grid/table layout with colored cells and divider lines
- Slide 8: A = horizontal timeline with circles | B = vertical step cards with numbered circles
- Slide 9: A = big stat number + image right | B = quote accent + image right

Rules:
- Each bullet must be concise (under 60 characters)
- Card titles: 2-4 words (or short stats like "50M", "98%" for slide 5B)
- stat_number: a short figure like "90%", "$2.5M", "10x", "47K"
- title_image_prompt: a vivid prompt (under 120 chars) for the title slide portrait image. Describe ONLY a photographic scene, object, or landscape. ABSOLUTELY NO text, letters, words, numbers, labels, signs, banners, captions, watermarks, logos, or writing of any kind. Never mention text-containing objects like books, screens, signs, or documents. Describe pure visual imagery only.
- image_prompt: a vivid prompt (under 120 chars) for AI image generation. Describe ONLY a photographic scene, object, or landscape. ABSOLUTELY NO text, letters, words, numbers, labels, signs, banners, captions, watermarks, logos, or writing of any kind. Never mention text-containing objects like books, screens, signs, or documents. Describe pure visual imagery only.
- All content must be relevant, factual, and professional
- Slide 10 must be a conclusion / call-to-action
- Do NOT wrap output in markdown code fences
- Mix designs A and B across slides for visual variety. Use at least 3 of each.
- For slide 2 always provide exactly 3 cards (no image needed)
- For slide 3A provide exactly 3 cards, for slide 3B provide exactly 4 cards
- For slide 5 always provide exactly 4 cards
- For slide 5B, card titles MUST be very short stat-like values (max 4 characters): "50M", "98%", "$2B", "10x". NEVER use full words or phrases as card titles for slide 5B.
- For slide 5A, card titles should be short (2-3 words max) so they fit neatly in cards.
- Choose slide designs wisely based on content and context to maximize DIVERSITY:
  * Use 2A (dot-badge rows) for overview/intro content; 2B (left-accent cards) for categorized info
  * Use 3A (3-column cards) for broad categories; 3B (4-column cards) for detailed breakdowns
  * Use 4A (image left) and 4B (image right) to alternate image placement
  * Use 5B (stat columns) ONLY when you have real numerical data; use 5A (grid) for concepts
  * Use 6A (image left + bullets right + banner) for data-heavy slides; 6B (image right + cards left + banner) for key takeaways
  * Use 7A (tinted comparison) for pros/cons; 7B (grid table) for detailed side-by-side comparisons
  * Use 8A (timeline) for sequential processes; 8B (step cards) for parallel steps
  * Use 9A (big stat + image) for impact numbers; 9B (quote + image) for testimonials/insights
- CRITICAL: Always use a healthy MIX of A and B designs. Never use all A or all B. Use at least 3 of each for visual variety. Each slide should feel different from the previous one.
- The presentation is for STUDENTS - make content engaging, clear, and easy to present. Use vivid language in bullets.
- "theme" must be one of: "indigo", "nature", "moss", "navy"
  Choose the theme that best fits the topic:
  indigo = design, art, fashion, luxury, creative industries, psychology
  nature = environment, health, biology, sustainability, food, agriculture
  moss   = education, social sciences, history, philosophy, literature
  navy   = business, technology, engineering, finance, law, science"""


def generate_slide_content(prompt):
    for attempt in range(2):
        try:
            r = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": _CONTENT_SYSTEM_PROMPT},
                    {"role": "user", "content": f"Create a presentation about: {prompt}"}
                ],
                temperature=0.7,
            )
            raw = r.choices[0].message.content.strip()
            # Strip markdown code fences if present
            if raw.startswith("```"):
                raw = raw.split("\n", 1)[1]
                if "```" in raw:
                    raw = raw.rsplit("```", 1)[0]
            # Try to extract JSON if there's extra text around it
            raw = raw.strip()
            if not raw.startswith("{"):
                start = raw.find("{")
                if start != -1:
                    raw = raw[start:]
            if not raw.endswith("}"):
                end = raw.rfind("}")
                if end != -1:
                    raw = raw[:end + 1]
            return json.loads(raw)
        except json.JSONDecodeError as e:
            print(f"[PresentAI] JSON parse error (attempt {attempt + 1}): {e}")
            if attempt == 0:
                continue
            return None
        except Exception as e:
            print(f"[PresentAI] generate_slide_content error (attempt {attempt + 1}): {e}")
            traceback.print_exc()
            if attempt == 0:
                continue
            return None
    return None


def generate_image(prompt, img_w=1024, img_h=1024):
    """Generate an image using HF FLUX.1-schnell. Returns BytesIO or None."""
    if not HF_API_KEY:
        print("[PresentAI] HF_API_KEY is empty — skipping image generation")
        return None

    # Try multiple API endpoints (router + direct inference)
    api_urls = [
        "https://router.huggingface.co/hf-inference/models/black-forest-labs/FLUX.1-schnell",
        "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-schnell",
    ]
    headers = {
        "Authorization": f"Bearer {HF_API_KEY}",
        "Content-Type": "application/json",
    }
    enhanced = (
        "Ultra sharp 4K professional photograph, crystal clear details, "
        "highly detailed, sharp focus, studio lighting, vibrant colors, "
        "no blurry elements, absolutely no text, no words, no letters, "
        "no numbers, no labels, no signs, no writing, no captions, "
        "no watermarks, no logos, text-free image only, " + prompt
    )
    payload = {
        "inputs": enhanced,
        "parameters": {
            "width": img_w,
            "height": img_h,
        },
    }

    for api_url in api_urls:
        for attempt in range(3):
            try:
                print(f"[PresentAI] Image attempt {attempt + 1} via {api_url.split('/')[2]}...")
                resp = requests.post(api_url, headers=headers, json=payload, timeout=120)
                if resp.status_code == 200:
                    ct = resp.headers.get("Content-Type", "")
                    if "image" in ct or len(resp.content) > 1000:
                        print(f"[PresentAI] Image generated ({len(resp.content)} bytes)")
                        return BytesIO(resp.content)
                    else:
                        print(f"[PresentAI] Unexpected response: {resp.text[:200]}")
                elif resp.status_code == 503:
                    try:
                        err = resp.json()
                        wait = min(int(err.get("estimated_time", 20)) + 1, 60)
                    except Exception:
                        wait = 20
                    print(f"[PresentAI] Model loading, waiting {wait}s...")
                    if attempt < 2:
                        time.sleep(wait)
                        continue
                else:
                    print(f"[PresentAI] Image API error {resp.status_code}: {resp.text[:200]}")
                    if resp.status_code in (401, 403):
                        print("[PresentAI] Auth error — check HF_API_KEY in .env")
                        return None
                    break  # Try next URL
            except requests.exceptions.Timeout:
                print(f"[PresentAI] Image request timed out (attempt {attempt + 1})")
            except Exception as e:
                print(f"[PresentAI] Image error: {e}")
        # If this URL failed all attempts, try next URL
    print("[PresentAI] All image generation attempts failed")
    return None


# ======================== UI COMPONENTS ========================

def _new_slide(prs, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color or SOFT_CREAM
    return slide


def add_slide_number(slide, num, total):
    box = slide.shapes.add_textbox(Inches(SW - MG - 1.0), Inches(SH - 0.55),
                                   Inches(1.0), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{num} / {total}"
    r.font.size = Pt(10)
    r.font.color.rgb = MEDIUM_GRAY


def add_accent_line(slide, left, top, width, color=None):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = color or WARM_CORAL
    line.line.fill.background()


def add_section_title(slide, text, left=None, top=None, width=None,
                      color=None, align=PP_ALIGN.LEFT):
    l = left if left is not None else Inches(MG)
    t = top if top is not None else Inches(MG)
    w = width if width is not None else Inches(CW)
    box = slide.shapes.add_textbox(l, t, w, Inches(0.90))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = color or CHARCOAL


def add_circle(slide, left, top, size, color):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c


def add_square_image(slide, image_bytes, left, top, size=IMG_SZ):
    """Place a square image at the given position."""
    if image_bytes:
        slide.shapes.add_picture(image_bytes, Inches(left), Inches(top),
                                  Inches(size), Inches(size))
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(left), Inches(top),
                                       Inches(size), Inches(size))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF0, 0xF1, 0xF5)
        rect.line.fill.background()
        tf = rect.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = Inches(size / 2 - 0.3)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "IMAGE"
        r.font.size = Pt(16)
        r.font.color.rgb = MEDIUM_GRAY


def add_bullet_text(slide, left, top, width, height, bullets,
                    bullet_color=None, font_size=13, text_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    b_color = bullet_color or ROYAL_BLUE
    t_color = text_color or SLATE

    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        p.line_spacing = Pt(max(font_size + 8, 20))
        rb = p.add_run()
        rb.text = "  \u2022  "
        rb.font.size = Pt(font_size)
        rb.font.bold = True
        rb.font.color.rgb = b_color
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(font_size)
        rt.font.color.rgb = t_color


def add_card(slide, left, top, width, height, bg_color=None, accent_color=None):
    """Rounded card with optional colored accent bar at top edge."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color or WHITE
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.75)
    card.adjustments[0] = CARD_RADIUS
    if accent_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left + Inches(0.08), top,
                                     width - Inches(0.16), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
    return card


def add_card_with_bullets(slide, left, top, width, height, title, bullets,
                          bg_color=None, accent_color=None,
                          title_color=None, bullet_color=None):
    add_card(slide, left, top, width, height, bg_color, accent_color)
    tbox = slide.shapes.add_textbox(left + Inches(0.20), top + Inches(0.15),
                                    width - Inches(0.40), Inches(0.50))
    tf = tbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = title_color or CHARCOAL
    add_bullet_text(slide, left + Inches(0.15), top + Inches(0.60),
                    width - Inches(0.30), height - Inches(0.70),
                    bullets,
                    bullet_color=bullet_color or accent_color or ROYAL_BLUE,
                    font_size=13)


def add_side_image(slide, image_bytes, side="left", width=SIDE_W, height=SIDE_H):
    """Place a full-height image flush to left or right edge."""
    left = 0.0 if side == "left" else SW - width
    if image_bytes:
        slide.shapes.add_picture(image_bytes, Inches(left), Inches(0),
                                  Inches(width), Inches(height))
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(left), Inches(0),
                                       Inches(width), Inches(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF0, 0xF1, 0xF5)
        rect.line.fill.background()
        tf = rect.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = Inches(height / 2 - 0.3)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "IMAGE"
        r.font.size = Pt(16)
        r.font.color.rgb = MEDIUM_GRAY


def add_dot_badge(slide, left, top, color=None):
    """Small colored dot badge (0.12" circle)."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top,
                                Inches(0.12), Inches(0.12))
    c.fill.solid()
    c.fill.fore_color.rgb = color or WARM_CORAL
    c.line.fill.background()
    return c


def add_card_left_accent(slide, left, top, width, height,
                          bg_color=None, accent_color=None):
    """Card with colored accent bar on left edge."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color or WHITE
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.75)
    card.adjustments[0] = CARD_RADIUS
    if accent_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top + Inches(0.08),
                                     Inches(0.07), height - Inches(0.16))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
    return card


# ======================== SLIDE LAYOUTS ========================

def slide_1_hero_title(prs, title, subtitle, image_bytes=None):
    """Slide 1 - Full-height image left + accent line, title, subtitle right."""
    slide = _new_slide(prs, SOFT_CREAM)
    # Full-height portrait image on left
    add_side_image(slide, image_bytes, side="left")
    # Text area on the right
    text_left = SIDE_W + 1.00
    text_w = SW - text_left - 0.80
    # Thick accent line above title (matches reference style)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(text_left), Inches(2.40),
                                   Inches(2.80), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = WARM_CORAL
    line.line.fill.background()
    # Title - large bold
    box = slide.shapes.add_textbox(Inches(text_left), Inches(2.70),
                                   Inches(text_w), Inches(2.60))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(52)
    r = p.add_run()
    r.text = title
    r.font.size = Pt(42)
    r.font.bold = True
    r.font.color.rgb = CHARCOAL
    # Subtitle - lighter, positioned lower with breathing room
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(text_left), Inches(5.70),
                                        Inches(text_w), Inches(0.80))
        stf = sbox.text_frame
        stf.clear()
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(20)
        sr.font.color.rgb = SLATE
    add_slide_number(slide, 1, TOTAL_SLIDES)


def slide_2a_dot_badge_rows(prs, data):
    """Slide 2A - Full-width dot-badge bullet rows (no image)."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:3]
    n = len(cards)
    row_h = 1.55
    row_gap = 0.20
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL]
    for i, c in enumerate(cards):
        y = CT + i * (row_h + row_gap)
        color = colors[i % 3]
        # Dot badge
        add_dot_badge(slide, Inches(MG + 0.02), Inches(y + 0.07), color)
        # Card heading
        hbox = slide.shapes.add_textbox(Inches(MG + 0.28), Inches(y),
                                         Inches(CW - 0.28), Inches(0.40))
        hf = hbox.text_frame
        hf.clear()
        hf.word_wrap = True
        hp = hf.paragraphs[0]
        hr = hp.add_run()
        hr.text = c["title"]
        hr.font.size = Pt(18)
        hr.font.bold = True
        hr.font.color.rgb = CHARCOAL
        # Bullets
        add_bullet_text(slide, Inches(MG + 0.28), Inches(y + 0.42),
                        Inches(CW - 0.28), Inches(row_h - 0.45),
                        c["bullets"], bullet_color=color, font_size=13)
    add_slide_number(slide, 2, TOTAL_SLIDES)


def slide_2b_left_accent_cards(prs, data):
    """Slide 2B - Full-width cards with left accent bars (no image)."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:3]
    n = len(cards)
    card_gap = 0.20
    card_h = min((SH - CT - MG - (n - 1) * card_gap) / n, 1.60)
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL]
    for i, c in enumerate(cards):
        y = CT + i * (card_h + card_gap)
        color = colors[i % 3]
        add_card_left_accent(slide, Inches(MG), Inches(y),
                             Inches(CW), Inches(card_h),
                             bg_color=WHITE, accent_color=color)
        # Card title
        hbox = slide.shapes.add_textbox(Inches(MG + 0.30), Inches(y + 0.12),
                                         Inches(CW - 0.50), Inches(0.40))
        hf = hbox.text_frame
        hf.clear()
        hf.word_wrap = True
        hp = hf.paragraphs[0]
        hr = hp.add_run()
        hr.text = c["title"]
        hr.font.size = Pt(17)
        hr.font.bold = True
        hr.font.color.rgb = CHARCOAL
        # Bullets
        add_bullet_text(slide, Inches(MG + 0.25), Inches(y + 0.50),
                        Inches(CW - 0.45), Inches(card_h - 0.60),
                        c["bullets"], bullet_color=color, font_size=13)
    add_slide_number(slide, 2, TOTAL_SLIDES)


def slide_3a_three_cards_row(prs, data):
    """Slide 3A - Three equal cards in a row with white fills."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:3]
    n = len(cards)
    card_w = (CW - (n - 1) * GAP) / n
    card_h = 5.00
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL]
    fills = [WHITE, WHITE, WHITE]
    for i, c in enumerate(cards):
        x = MG + i * (card_w + GAP)
        add_card_with_bullets(slide, Inches(x), Inches(CT),
                              Inches(card_w), Inches(card_h),
                              c["title"], c["bullets"],
                              bg_color=fills[i], accent_color=colors[i % 3])
    add_slide_number(slide, 3, TOTAL_SLIDES)


def slide_3b_four_cards_row(prs, data):
    """Slide 3B - Four equal cards in a row with tinted fills."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:4]
    n = min(len(cards), 4)
    card_w = (CW - (n - 1) * GAP) / n
    card_h = 5.00
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL, ROYAL_BLUE]
    fills = [LIGHT_BLUE_TINT, LIGHT_CORAL_TINT, LIGHT_BLUE_TINT, LIGHT_CORAL_TINT]
    for i, c in enumerate(cards[:n]):
        x = MG + i * (card_w + GAP)
        add_card_with_bullets(slide, Inches(x), Inches(CT),
                              Inches(card_w), Inches(card_h),
                              c["title"], c["bullets"],
                              bg_color=fills[i], accent_color=colors[i % 4])
    add_slide_number(slide, 3, TOTAL_SLIDES)


def slide_4a_image_left_bullets(prs, data, image_bytes=None):
    """Slide 4A - Square image left + bullet list right."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    add_square_image(slide, image_bytes, MG, CT)
    txt_x = MG + IMG_SZ + GAP
    add_bullet_text(slide, Inches(txt_x), Inches(CT),
                    Inches(TXT_BW), Inches(5.00),
                    data["bullets"], bullet_color=WARM_CORAL, font_size=15)
    add_slide_number(slide, 4, TOTAL_SLIDES)


def slide_4b_bullets_image_right(prs, data, image_bytes=None):
    """Slide 4B - Bullet list left + square image right."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    add_square_image(slide, image_bytes, IMG_RX, IMG_TOP)
    add_bullet_text(slide, Inches(MG), Inches(CT),
                    Inches(TXT_BW), Inches(5.00),
                    data["bullets"], bullet_color=WARM_CORAL, font_size=15)
    add_slide_number(slide, 4, TOTAL_SLIDES)


def slide_5a_grid_badges(prs, data):
    """Slide 5A - 2x2 grid cards with numbered circle badges."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:4]
    card_w = (CW - GAP) / 2
    v_gap = 0.24
    card_h = (5.00 - v_gap) / 2
    positions = [
        (MG, CT),
        (MG + card_w + GAP, CT),
        (MG, CT + card_h + v_gap),
        (MG + card_w + GAP, CT + card_h + v_gap),
    ]
    accent_colors = [TEAL, TEAL, TEAL, TEAL]
    fills = [WHITE, LIGHT_BLUE_TINT, LIGHT_BLUE_TINT, WHITE]
    for i, c in enumerate(cards):
        lx, ly = positions[i]
        add_card_with_bullets(slide, Inches(lx), Inches(ly),
                              Inches(card_w), Inches(card_h),
                              c["title"], c["bullets"],
                              bg_color=fills[i], accent_color=accent_colors[i])
        badge = add_circle(slide,
                           Inches(lx + card_w - 0.55),
                           Inches(ly + 0.10),
                           Inches(0.40), TEAL)
        ctf = badge.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
    add_slide_number(slide, 5, TOTAL_SLIDES)


def slide_5b_stat_columns(prs, data):
    """Slide 5B - Four stat cards with highlighted numbers."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    cards = data["cards"][:4]
    n = len(cards)
    card_w = (CW - (n - 1) * GAP) / n
    card_h = 5.00
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL, ROYAL_BLUE]
    for i, c in enumerate(cards):
        x = MG + i * (card_w + GAP)
        color = colors[i % 4]
        add_card(slide, Inches(x), Inches(CT),
                 Inches(card_w), Inches(card_h),
                 bg_color=WHITE, accent_color=color)
        title_text = c["title"]
        title_sz = 20 if len(title_text) <= 25 else 16
        nbox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(CT + 0.20),
                                         Inches(card_w - 0.30), Inches(0.65))
        nf = nbox.text_frame
        nf.clear()
        nf.word_wrap = True
        np_ = nf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = title_text
        nr.font.size = Pt(title_sz)
        nr.font.bold = True
        nr.font.color.rgb = color
        line_w = min(0.8, card_w * 0.35)
        line_y = CT + 0.95
        add_accent_line(slide, Inches(x + (card_w - line_w) / 2),
                        Inches(line_y), Inches(line_w), color)
        add_bullet_text(slide, Inches(x + 0.10), Inches(line_y + 0.18),
                        Inches(card_w - 0.20), Inches(card_h - (line_y - CT) - 0.33),
                        c["bullets"],
                        bullet_color=color, font_size=12)
    add_slide_number(slide, 5, TOTAL_SLIDES)


def slide_6a_bullets_banner(prs, data, image_bytes=None):
    """Slide 6A - Image left + bullet list right + dark stat banner at bottom."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    banner_h = 1.10
    banner_top = SH - banner_h
    # Image on left
    img_sz = 4.50
    add_square_image(slide, image_bytes, MG, CT, size=img_sz)
    # Bullets beside image
    txt_x = MG + img_sz + GAP
    txt_w = SW - txt_x - MG
    bullet_h = banner_top - CT - 0.20
    add_bullet_text(slide, Inches(txt_x), Inches(CT),
                    Inches(txt_w), Inches(bullet_h),
                    data["bullets"], bullet_color=ROYAL_BLUE, font_size=15)
    # Dark stat banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(banner_top),
                                    Inches(SW), Inches(banner_h))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DEEP_NAVY
    banner.line.fill.background()
    tf = banner.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(MG)
    tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = data.get("stat_number", "")
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = WARM_CORAL
    r2 = p.add_run()
    r2.text = "  " + data.get("stat_label", "")
    r2.font.size = Pt(18)
    r2.font.color.rgb = WHITE
    add_slide_number(slide, 6, TOTAL_SLIDES)


def slide_6b_cards_banner(prs, data, image_bytes=None):
    """Slide 6B - Image right + stacked cards left + dark stat banner at bottom."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    banner_h = 1.10
    banner_top = SH - banner_h
    # Image on right
    img_sz = 4.50
    add_square_image(slide, image_bytes, IMG_RX, IMG_TOP, size=img_sz)
    # Cards on left beside image
    card_w = IMG_RX - MG - GAP
    bullets = data["bullets"]
    n = min(len(bullets), 5)
    avail_h = banner_top - CT - 0.20
    card_gap = 0.15
    card_h = min((avail_h - (n - 1) * card_gap) / n, 0.80)
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL, ROYAL_BLUE, WARM_CORAL]
    for i, bullet in enumerate(bullets[:n]):
        y = CT + i * (card_h + card_gap)
        add_card(slide, Inches(MG), Inches(y),
                 Inches(card_w), Inches(card_h),
                 bg_color=WHITE, accent_color=colors[i % 5])
        tbox = slide.shapes.add_textbox(Inches(MG + 0.25), Inches(y + 0.12),
                                        Inches(card_w - 0.50), Inches(card_h - 0.24))
        ttf = tbox.text_frame
        ttf.clear()
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        rb = tp.add_run()
        rb.text = "  \u2022  "
        rb.font.size = Pt(15)
        rb.font.bold = True
        rb.font.color.rgb = colors[i % 5]
        rt = tp.add_run()
        rt.text = bullet
        rt.font.size = Pt(15)
        rt.font.color.rgb = SLATE
    # Dark stat banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(banner_top),
                                    Inches(SW), Inches(banner_h))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DEEP_NAVY
    banner.line.fill.background()
    tf = banner.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(MG)
    tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = data.get("stat_number", "")
    r1.font.size = Pt(28)
    r1.font.bold = True
    r1.font.color.rgb = WARM_CORAL
    r2 = p.add_run()
    r2.text = "  " + data.get("stat_label", "")
    r2.font.size = Pt(18)
    r2.font.color.rgb = WHITE
    add_slide_number(slide, 6, TOTAL_SLIDES)


def slide_7a_two_cards(prs, data):
    """Slide 7A - Two large side-by-side cards with tinted fills."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    card_w = (CW - GAP) / 2
    card_h = 5.00
    add_card_with_bullets(slide, Inches(MG), Inches(CT),
                          Inches(card_w), Inches(card_h),
                          data["left"]["title"], data["left"]["bullets"],
                          bg_color=LIGHT_BLUE_TINT, accent_color=ROYAL_BLUE,
                          bullet_color=ROYAL_BLUE)
    add_card_with_bullets(slide, Inches(MG + card_w + GAP), Inches(CT),
                          Inches(card_w), Inches(card_h),
                          data["right"]["title"], data["right"]["bullets"],
                          bg_color=LIGHT_CORAL_TINT, accent_color=WARM_CORAL,
                          bullet_color=WARM_CORAL)
    add_slide_number(slide, 7, TOTAL_SLIDES)


def slide_7b_grid_table(prs, data):
    """Slide 7B - Grid/table layout with colored cells and dividers."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    left_data = data["left"]
    right_data = data["right"]
    card_w = (CW - GAP) / 2
    card_h = 5.00
    # Left card - tinted fill
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(MG), Inches(CT),
                                        Inches(card_w), Inches(card_h))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = LIGHT_BLUE_TINT
    left_card.line.color.rgb = CARD_BORDER
    left_card.line.width = Pt(0.75)
    left_card.adjustments[0] = CARD_RADIUS
    # Left title
    ltbox = slide.shapes.add_textbox(Inches(MG + 0.25), Inches(CT + 0.20),
                                      Inches(card_w - 0.50), Inches(0.50))
    ltf = ltbox.text_frame
    ltf.clear()
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_data["title"]
    lr.font.size = Pt(18)
    lr.font.bold = True
    lr.font.color.rgb = ROYAL_BLUE
    # Left bullets as rows with dividers
    bullets_l = left_data["bullets"]
    row_start = CT + 0.80
    row_h = (card_h - 1.00) / max(len(bullets_l), 1)
    for j, b in enumerate(bullets_l):
        y = row_start + j * row_h
        if j > 0:
            divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(MG + 0.15), Inches(y),
                                              Inches(card_w - 0.30), Pt(0.75))
            divider.fill.solid()
            divider.fill.fore_color.rgb = CARD_BORDER
            divider.line.fill.background()
        bx = slide.shapes.add_textbox(Inches(MG + 0.25), Inches(y + 0.06),
                                       Inches(card_w - 0.50), Inches(row_h - 0.08))
        bf = bx.text_frame
        bf.clear()
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        br1 = bp.add_run()
        br1.text = "  \u2022  "
        br1.font.size = Pt(13)
        br1.font.bold = True
        br1.font.color.rgb = ROYAL_BLUE
        br2 = bp.add_run()
        br2.text = b
        br2.font.size = Pt(13)
        br2.font.color.rgb = SLATE
    # Right card - different tint
    rx = MG + card_w + GAP
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(rx), Inches(CT),
                                         Inches(card_w), Inches(card_h))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_CORAL_TINT
    right_card.line.color.rgb = CARD_BORDER
    right_card.line.width = Pt(0.75)
    right_card.adjustments[0] = CARD_RADIUS
    # Right title
    rtbox = slide.shapes.add_textbox(Inches(rx + 0.25), Inches(CT + 0.20),
                                      Inches(card_w - 0.50), Inches(0.50))
    rtf = rtbox.text_frame
    rtf.clear()
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_data["title"]
    rr.font.size = Pt(18)
    rr.font.bold = True
    rr.font.color.rgb = WARM_CORAL
    # Right bullets as rows with dividers
    bullets_r = right_data["bullets"]
    for j, b in enumerate(bullets_r):
        y = row_start + j * row_h
        if j > 0:
            divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(rx + 0.15), Inches(y),
                                              Inches(card_w - 0.30), Pt(0.75))
            divider.fill.solid()
            divider.fill.fore_color.rgb = CARD_BORDER
            divider.line.fill.background()
        bx = slide.shapes.add_textbox(Inches(rx + 0.25), Inches(y + 0.06),
                                       Inches(card_w - 0.50), Inches(row_h - 0.08))
        bf = bx.text_frame
        bf.clear()
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        br1 = bp.add_run()
        br1.text = "  \u2022  "
        br1.font.size = Pt(13)
        br1.font.bold = True
        br1.font.color.rgb = WARM_CORAL
        br2 = bp.add_run()
        br2.text = b
        br2.font.size = Pt(13)
        br2.font.color.rgb = SLATE
    add_slide_number(slide, 7, TOTAL_SLIDES)


def slide_8a_timeline(prs, data):
    """Slide 8A - Horizontal timeline with numbered circles."""
    slide = _new_slide(prs, LIGHT_BLUE_TINT)
    add_section_title(slide, data["title"])
    steps = data["steps"][:4]
    n = len(steps)
    step_w = CW / n
    line_y = 2.87
    conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(MG + 0.40), Inches(line_y),
                                  Inches(CW - 0.80), Inches(0.06))
    conn.fill.solid()
    conn.fill.fore_color.rgb = MEDIUM_GRAY
    conn.line.fill.background()
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL, ROYAL_BLUE]
    for i, step in enumerate(steps):
        col_left = MG + i * step_w
        cx = col_left + step_w / 2 - 0.25
        c = add_circle(slide, Inches(cx), Inches(line_y - 0.10),
                        Inches(0.50), colors[i % 4])
        ctf = c.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        cr.font.size = Pt(16)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
        tbox = slide.shapes.add_textbox(Inches(col_left), Inches(3.40),
                                        Inches(step_w), Inches(0.50))
        ttf = tbox.text_frame
        ttf.clear()
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = step["title"]
        tr.font.size = Pt(14)
        tr.font.bold = True
        tr.font.color.rgb = CHARCOAL
        add_bullet_text(slide, Inches(col_left), Inches(4.00),
                        Inches(step_w), Inches(2.80),
                        step["bullets"], bullet_color=colors[i % 4],
                        font_size=12)
    add_slide_number(slide, 8, TOTAL_SLIDES)


def slide_8b_step_cards(prs, data):
    """Slide 8B - Vertical step cards with numbered circles."""
    slide = _new_slide(prs, LIGHT_BLUE_TINT)
    add_section_title(slide, data["title"])
    steps = data["steps"][:4]
    n = len(steps)
    card_gap = 0.20
    card_h = min((5.00 - (n - 1) * card_gap) / n, 1.10)
    colors = [ROYAL_BLUE, WARM_CORAL, TEAL, ROYAL_BLUE]
    for i, step in enumerate(steps):
        y = CT + i * (card_h + card_gap)
        add_card(slide, Inches(MG + 0.65), Inches(y),
                 Inches(CW - 0.65), Inches(card_h),
                 bg_color=WHITE, accent_color=colors[i % 4])
        tbox = slide.shapes.add_textbox(Inches(MG + 0.85), Inches(y + 0.10),
                                        Inches(CW - 1.10), Inches(0.35))
        tf = tbox.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step["title"]
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL
        btext = " \u00b7 ".join(step["bullets"])
        bbox = slide.shapes.add_textbox(Inches(MG + 0.85), Inches(y + 0.42),
                                        Inches(CW - 1.10), Inches(card_h - 0.52))
        bf = bbox.text_frame
        bf.clear()
        bf.word_wrap = True
        bp = bf.paragraphs[0]
        br = bp.add_run()
        br.text = btext
        br.font.size = Pt(12)
        br.font.color.rgb = SLATE
        c = add_circle(slide, Inches(MG), Inches(y + card_h / 2 - 0.25),
                        Inches(0.50), colors[i % 4])
        ctf = c.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        cr.font.size = Pt(16)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
    add_slide_number(slide, 8, TOTAL_SLIDES)


def slide_9a_stat_image(prs, data, image_bytes=None):
    """Slide 9A - Big stat number + bullets left, square image right."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    add_square_image(slide, image_bytes, IMG_RX, IMG_TOP)
    text_w = TXT_BW
    sbox = slide.shapes.add_textbox(Inches(MG), Inches(CT),
                                    Inches(text_w), Inches(1.50))
    stf = sbox.text_frame
    stf.clear()
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = data.get("stat_number", "")
    sr.font.size = Pt(64)
    sr.font.bold = True
    sr.font.color.rgb = WARM_CORAL
    sp2 = stf.add_paragraph()
    sr2 = sp2.add_run()
    sr2.text = data.get("stat_label", "")
    sr2.font.size = Pt(18)
    sr2.font.color.rgb = SLATE
    add_bullet_text(slide, Inches(MG), Inches(CT + 2.20),
                    Inches(text_w), Inches(2.80),
                    data["bullets"], bullet_color=TEAL, font_size=14)
    add_slide_number(slide, 9, TOTAL_SLIDES)


def slide_9b_quote_image(prs, data, image_bytes=None):
    """Slide 9B - Quote with vertical accent line + square image right."""
    slide = _new_slide(prs)
    add_section_title(slide, data["title"])
    add_square_image(slide, image_bytes, IMG_RX, IMG_TOP)
    text_w = TXT_BW
    quote = f"{data.get('stat_number', '')} {data.get('stat_label', '')}"
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(MG), Inches(CT + 0.10),
                                   Pt(5), Inches(1.80))
    line.fill.solid()
    line.fill.fore_color.rgb = WARM_CORAL
    line.line.fill.background()
    qbox = slide.shapes.add_textbox(Inches(MG + 0.30), Inches(CT),
                                     Inches(text_w - 0.30), Inches(2.00))
    qf = qbox.text_frame
    qf.clear()
    qf.word_wrap = True
    qp = qf.paragraphs[0]
    qr = qp.add_run()
    qr.text = quote
    qr.font.size = Pt(28)
    qr.font.bold = True
    qr.font.color.rgb = CHARCOAL
    add_bullet_text(slide, Inches(MG), Inches(CT + 2.50),
                    Inches(text_w), Inches(2.50),
                    data["bullets"], bullet_color=TEAL, font_size=14)
    add_slide_number(slide, 9, TOTAL_SLIDES)


def slide_10_closing_cta(prs, data):
    """Slide 10 - Dark closing with action cards."""
    slide = _new_slide(prs, DEEP_NAVY)
    title_x = 2.10
    title_w = SW - 2 * title_x
    box = slide.shapes.add_textbox(Inches(title_x), Inches(MG),
                                   Inches(title_w), Inches(1.20))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = data["title"]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    div_w = 3.00
    add_accent_line(slide, Inches((SW - div_w) / 2), Inches(1.85),
                    Inches(div_w), WARM_CORAL)
    card_w = (CW - 2 * GAP) / 3
    card_top = 2.30
    card_h = 4.50
    for i, c in enumerate(data["cards"][:3]):
        x = MG + i * (card_w + GAP)
        add_card_with_bullets(slide, Inches(x), Inches(card_top),
                              Inches(card_w), Inches(card_h),
                              c["title"], c["bullets"],
                              bg_color=LIGHT_CORAL_TINT,
                              accent_color=WARM_CORAL,
                              bullet_color=WARM_CORAL)
    add_slide_number(slide, 10, TOTAL_SLIDES)


# ======================== GENERATE PPT ========================

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    if "user" not in session:
        return redirect("/login")

    user_prompt = request.form["user_prompt"]

    content = generate_slide_content(user_prompt)
    if content is None:
        return jsonify({"error": "Failed to generate presentation content. Please try again."}), 500

    raw_title = content.get("title", user_prompt)
    subtitle = content.get("subtitle", "")
    filename = safe_filename(raw_title)

    apply_theme(content.get("theme", "indigo"))

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    slides = {s["slide"]: s for s in content.get("slides", [])}
    designs = content.get("slide_designs", {})

    # Generate 4 images: slide 1 (portrait), slides 4, 6, 9 (square)
    images = {}
    title_img_prompt = content.get("title_image_prompt", "")
    if title_img_prompt:
        images[1] = generate_image(title_img_prompt, img_w=768, img_h=1152)
    for sn in [4, 6, 9]:
        prompt = slides.get(sn, {}).get("image_prompt", "")
        if prompt:
            images[sn] = generate_image(prompt)

    slide_1_hero_title(prs, raw_title, subtitle, images.get(1))

    if designs.get("2", "A") == "B":
        slide_2b_left_accent_cards(prs, slides[2])
    else:
        slide_2a_dot_badge_rows(prs, slides[2])

    if designs.get("3", "A") == "B":
        slide_3b_four_cards_row(prs, slides[3])
    else:
        slide_3a_three_cards_row(prs, slides[3])

    if designs.get("4", "A") == "B":
        slide_4b_bullets_image_right(prs, slides[4], images.get(4))
    else:
        slide_4a_image_left_bullets(prs, slides[4], images.get(4))

    if designs.get("5", "A") == "B":
        slide_5b_stat_columns(prs, slides[5])
    else:
        slide_5a_grid_badges(prs, slides[5])

    if designs.get("6", "A") == "B":
        slide_6b_cards_banner(prs, slides[6], images.get(6))
    else:
        slide_6a_bullets_banner(prs, slides[6], images.get(6))

    if designs.get("7", "A") == "B":
        slide_7b_grid_table(prs, slides[7])
    else:
        slide_7a_two_cards(prs, slides[7])

    if designs.get("8", "A") == "B":
        slide_8b_step_cards(prs, slides[8])
    else:
        slide_8a_timeline(prs, slides[8])

    if designs.get("9", "A") == "B":
        slide_9b_quote_image(prs, slides[9], images.get(9))
    else:
        slide_9a_stat_image(prs, slides[9], images.get(9))

    slide_10_closing_cta(prs, slides[10])

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Save to history
    file_bytes = buf.getvalue()
    try:
        db, cursor = get_db()
        cursor.execute(
            "INSERT INTO presentations (user_email, title, filename, file_data) VALUES (%s, %s, %s, %s)",
            (session.get("email", ""), raw_title, f"{filename}.pptx", file_bytes),
        )
        db.commit()
    except Exception:
        pass  # Don't block download if history save fails

    buf.seek(0)
    return send_file(
        buf,
        as_attachment=True,
        download_name=f"{filename}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# ------------------ Auth Pages ------------------
@app.route("/")
def home():
    logged_in = "user" in session
    return render_template("home.html", logged_in=logged_in)

@app.route("/login")
def login_page():
    return render_template("login.html")

@app.route("/login", methods=["POST"])
@db_required
def login():
    db, cursor = get_db()
    email = request.form["email"]
    password = request.form["password"]
    cursor.execute("SELECT * FROM users WHERE email=%s", (email,))
    user = cursor.fetchone()
    if user:
        stored_password = user[4]  # password_hash column
        # Support both hashed and legacy plain-text passwords
        if stored_password and (
            check_password_hash(stored_password, password)
            or stored_password == password
        ):
            session["user"] = user[2] or email  # first_name or email
            session["email"] = email
            # Restore language preference
            lang = user[5] or "en"
            lang_map = {"en": "English", "es": "Spanish", "fr": "French", "de": "German",
                        "it": "Italian", "pt": "Portuguese", "ru": "Russian", "zh": "Chinese",
                        "ja": "Japanese", "ko": "Korean", "ar": "Arabic", "hi": "Hindi", "uz": "Uzbek"}
            session["language"] = lang_map.get(lang, "English")
            return redirect("/dashboard")
    return render_template("login.html", error="Invalid email or password")

@app.route("/register")
def register_page():
    return render_template("register.html")

@app.route("/register", methods=["POST"])
@db_required
def register():
    db, cursor = get_db()
    email = request.form["email"]
    first_name = request.form["first_name"]
    last_name = request.form["last_name"]
    language = request.form.get("language", "English")
    password = request.form["password"]

    # Check if email already exists
    cursor.execute("SELECT id FROM users WHERE email=%s", (email,))
    if cursor.fetchone():
        return render_template("register.html", error="An account with this email already exists")

    hashed_password = generate_password_hash(password)
    cursor.execute(
        "INSERT INTO users (email, first_name, last_name, language, password_hash) VALUES (%s, %s, %s, %s, %s)",
        (email, first_name, last_name, language, hashed_password),
    )
    db.commit()

    session["user"] = first_name
    session["email"] = email
    session["language"] = language
    return redirect("/dashboard")

# ------------------ Google OAuth ------------------
@app.route("/auth/google")
def google_login():
    redirect_uri = url_for("google_callback", _external=True)
    return google.authorize_redirect(redirect_uri)

@app.route("/auth/google/callback")
@db_required
def google_callback():
    try:
        token = google.authorize_access_token()
    except Exception:
        # State mismatch or expired — restart the OAuth flow
        return redirect("/auth/google")
    user_info = token.get("userinfo")
    if not user_info:
        return redirect("/login")

    email = user_info["email"]
    first_name = user_info.get("given_name", "")
    last_name = user_info.get("family_name", "")
    google_id = user_info["sub"]

    db, cursor = get_db()
    cursor.execute("SELECT * FROM users WHERE email=%s", (email,))
    existing = cursor.fetchone()

    if existing:
        # Update google_id if not set
        if not existing[8]:
            cursor.execute("UPDATE users SET google_id=%s WHERE email=%s", (google_id, email))
            db.commit()
        session["user"] = existing[2] or email  # first_name or email
    else:
        cursor.execute(
            "INSERT INTO users (email, first_name, last_name, password_hash, google_id) VALUES (%s, %s, %s, %s, %s)",
            (email, first_name, last_name, "", google_id),
        )
        db.commit()
        session["user"] = first_name or email

    session["email"] = email
    # Set language from DB for Google users
    if existing and existing[5]:
        lang = existing[5]
        lang_map = {"en": "English", "es": "Spanish", "fr": "French", "de": "German",
                    "it": "Italian", "pt": "Portuguese", "ru": "Russian", "zh": "Chinese",
                    "ja": "Japanese", "ko": "Korean", "ar": "Arabic", "hi": "Hindi", "uz": "Uzbek"}
        session["language"] = lang_map.get(lang, "English")
    else:
        session["language"] = "English"
    return redirect("/dashboard")

@app.route("/dashboard")
def dashboard():
    if "user" not in session:
        return redirect("/login")
    t = get_translations()
    return render_template("dashboard.html", user=session["user"], t=t)

@app.route("/profile")
@db_required
def profile():
    if "user" not in session:
        return redirect("/login")
    db, cursor = get_db()
    cursor.execute("SELECT * FROM users WHERE email=%s", (session.get("email", ""),))
    user_data = cursor.fetchone()
    t = get_translations()
    return render_template("profile.html", user=session["user"], user_data=user_data, t=t)

@app.route("/history")
@db_required
def history():
    if "user" not in session:
        return redirect("/login")
    db, cursor = get_db()
    cursor.execute(
        "SELECT id, title, filename, created_at FROM presentations WHERE user_email=%s ORDER BY created_at DESC",
        (session.get("email", ""),),
    )
    presentations = cursor.fetchall()
    t = get_translations()
    return render_template("history.html", user=session["user"], presentations=presentations, t=t)

@app.route("/download/<int:pres_id>")
@db_required
def download_presentation(pres_id):
    if "user" not in session:
        return redirect("/login")
    db, cursor = get_db()
    cursor.execute(
        "SELECT filename, file_data FROM presentations WHERE id=%s AND user_email=%s",
        (pres_id, session.get("email", "")),
    )
    row = cursor.fetchone()
    if not row:
        return redirect("/history")
    buf = BytesIO()
    prs.save(buf)

    # Save to history using getvalue() (doesn't move the pointer)
    file_bytes = buf.getvalue()
    try:
        db, cursor = get_db()
        cursor.execute(
            "INSERT INTO presentations (user_email, title, filename, file_data) VALUES (%s, %s, %s, %s)",
            (session.get("email", ""), raw_title, f"{filename}.pptx", file_bytes),
        )
        db.commit()
    except Exception:
        pass  # Don't block download if history save fails

    # Create a FRESH BytesIO from the saved bytes for send_file
    download_buf = BytesIO(file_bytes)
    return send_file(
        download_buf,
        as_attachment=True,
        download_name=f"{filename}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.route("/delete-account", methods=["POST"])
@db_required
def delete_account():
    if "user" not in session:
        return redirect("/login")
    email = session.get("email", "")
    try:
        db, cursor = get_db()
        cursor.execute("DELETE FROM presentations WHERE user_email=%s", (email,))
        cursor.execute("DELETE FROM users WHERE email=%s", (email,))
        db.commit()
    except Exception:
        pass
    session.pop("user", None)
    session.pop("email", None)
    session.pop("language", None)
    return redirect("/")

@app.route("/logout")
def logout():
    session.pop("user", None)
    session.pop("email", None)
    session.pop("language", None)
    return redirect("/")

# ------------------ Run ------------------
if __name__ == "__main__":
    port = int(os.getenv("PORT", 8000))
    debug = os.getenv("FLASK_DEBUG", "false").lower() == "true"
    app.run(debug=debug, host="0.0.0.0", port=port)
